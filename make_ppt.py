#!/usr/bin/env python3
"""Generate antirev architecture presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ─────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x1B, 0x1B, 0x2F)   # dark navy
BG_CARD   = RGBColor(0x25, 0x25, 0x3D)   # card bg
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
ORANGE    = RGBColor(0xFF, 0xA6, 0x2B)   # orange accent
GREEN     = RGBColor(0x06, 0xD6, 0xA0)   # green
RED       = RGBColor(0xEF, 0x47, 0x6F)   # red/pink
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xA0, 0xA0, 0xB0)
LIGHT     = RGBColor(0xE0, 0xE0, 0xE8)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return p


def add_para(tf, text, size=16, color=LIGHT, bold=False, space_before=Pt(4),
             space_after=Pt(2), level=0, alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    p.alignment = alignment
    return p


def add_bullet(tf, text, size=15, color=LIGHT, level=0):
    return add_para(tf, text, size=size, color=color, level=level,
                    space_before=Pt(3), space_after=Pt(1))


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD,
                     border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.7), Inches(2), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    tb = add_textbox(slide, Inches(1.5), Inches(1.0), Inches(7), Inches(1.6))
    set_text(tb.text_frame, title, size=40, color=WHITE, bold=True)

    tb = add_textbox(slide, Inches(1.5), Inches(2.9), Inches(7), Inches(1.2))
    set_text(tb.text_frame, subtitle, size=20, color=GRAY)


def section_slide(prs, number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Big number
    tb = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(1.5), Inches(2))
    set_text(tb.text_frame, "%02d" % number, size=72, color=ACCENT, bold=True)

    # Title
    tb = add_textbox(slide, Inches(2.5), Inches(1.5), Inches(6.5), Inches(1.2))
    set_text(tb.text_frame, title, size=36, color=WHITE, bold=True)

    if subtitle:
        tb = add_textbox(slide, Inches(2.5), Inches(2.5), Inches(6.5), Inches(1.0))
        set_text(tb.text_frame, subtitle, size=18, color=GRAY)

    return slide


def content_slide(prs, title, bullets, subbullets=None):
    """Simple bullet-list slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()

    tb = add_textbox(slide, Inches(0.6), Inches(0.1), Inches(8.5), Inches(0.7))
    set_text(tb.text_frame, title, size=26, color=ACCENT, bold=True)

    # Body
    tb = add_textbox(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(6.0))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for i, b in enumerate(bullets):
        if first:
            set_text(tf, b, size=15, color=LIGHT)
            first = False
        else:
            add_bullet(tf, b, size=15, color=LIGHT, level=0)
        # Add sub-bullets if any
        if subbullets and i in subbullets:
            for sb in subbullets[i]:
                add_bullet(tf, sb, size=13, color=GRAY, level=1)

    return slide


def two_col_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()

    tb = add_textbox(slide, Inches(0.6), Inches(0.1), Inches(8.5), Inches(0.7))
    set_text(tb.text_frame, title, size=26, color=ACCENT, bold=True)

    # Left column
    add_rounded_rect(slide, Inches(0.4), Inches(1.1), Inches(4.3), Inches(5.8),
                     border_color=ACCENT)
    tb = add_textbox(slide, Inches(0.6), Inches(1.2), Inches(4.0), Inches(0.5))
    set_text(tb.text_frame, left_title, size=18, color=ORANGE, bold=True)

    tb = add_textbox(slide, Inches(0.6), Inches(1.7), Inches(4.0), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(left_bullets):
        if i == 0:
            set_text(tf, b, size=14, color=LIGHT)
        else:
            add_bullet(tf, b, size=14, color=LIGHT)

    # Right column
    add_rounded_rect(slide, Inches(5.2), Inches(1.1), Inches(4.3), Inches(5.8),
                     border_color=GREEN)
    tb = add_textbox(slide, Inches(5.4), Inches(1.2), Inches(4.0), Inches(0.5))
    set_text(tb.text_frame, right_title, size=18, color=GREEN, bold=True)

    tb = add_textbox(slide, Inches(5.4), Inches(1.7), Inches(4.0), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(right_bullets):
        if i == 0:
            set_text(tf, b, size=14, color=LIGHT)
        else:
            add_bullet(tf, b, size=14, color=LIGHT)

    return slide


def diagram_flow(prs, title, steps, note=""):
    """Horizontal flow diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()

    tb = add_textbox(slide, Inches(0.6), Inches(0.1), Inches(8.5), Inches(0.7))
    set_text(tb.text_frame, title, size=26, color=ACCENT, bold=True)

    n = len(steps)
    box_w = min(Inches(1.8), Inches(8.5) / n)
    gap = Inches(0.15)
    total_w = n * box_w + (n - 1) * gap
    start_x = (Inches(10) - total_w) // 2
    top_y = Inches(2.0)
    box_h = Inches(2.2)

    colors = [ACCENT, GREEN, ORANGE, RED, ACCENT2, GREEN]

    for i, (label, desc) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        c = colors[i % len(colors)]
        add_rounded_rect(slide, x, top_y, box_w, box_h, border_color=c)

        # Step number
        tb = add_textbox(slide, x, top_y + Inches(0.1), box_w, Inches(0.35))
        set_text(tb.text_frame, label, size=13, color=c, bold=True,
                 alignment=PP_ALIGN.CENTER)

        # Description
        tb = add_textbox(slide, x + Inches(0.08), top_y + Inches(0.45),
                         box_w - Inches(0.16), box_h - Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        set_text(tf, desc, size=11, color=LIGHT)

        # Arrow between boxes
        if i < n - 1:
            ax = x + box_w + Inches(0.02)
            ay = top_y + box_h // 2
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, ax, ay - Inches(0.12),
                gap - Inches(0.04), Inches(0.24))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()

    if note:
        tb = add_textbox(slide, Inches(0.6), Inches(5.0), Inches(8.8), Inches(1.5))
        set_text(tb.text_frame, note, size=13, color=GRAY)

    return slide


# ═══════════════════════════════════════════════════════════════════════
#  BUILD THE PRESENTATION
# ═══════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ── Slide 1: Title ──────────────────────────────────────────────────
title_slide(prs,
    "AntiRev",
    "Binary Protection System\n"
    "Encrypt executables & shared libraries, run from memory\n"
    "Prevent reverse engineering at the deployment layer")

# ── Slide 2: Problem & Goals ────────────────────────────────────────
content_slide(prs, "Problem Statement & Goals", [
    "Business software deployed on customer-controlled hardware",
    "100+ executables, 550+ shared libraries, 1000+ Python scripts",
    "Threat: reverse engineering, binary analysis, IP extraction",
    "",
    "Goals:",
    "Plaintext binaries never exist on disk at runtime",
    "Transparent to the application -- no source code changes",
    "Support complex dependency graphs (DT_NEEDED, dlopen, Python ctypes)",
    "Scale to 550+ encrypted libs across 100+ concurrent processes",
    "Preserve original symbol resolution order and ABI semantics",
])

# ── Slide 3: Architecture Overview ──────────────────────────────────
section_slide(prs, 1, "Architecture Overview",
              "Six core components working together")

content_slide(prs, "Core Components", [
    "Stub  --  C launcher: decrypts, creates memfds, fexecve",
    "exe_shim  --  LD_PRELOAD: hides memfd paths from /proc/self/exe",
    "dlopen_shim  --  LD_PRELOAD: redirects dlopen() to encrypted libs",
    "Daemon (.antirev-libd)  --  Decrypt once, serve many via SCM_RIGHTS",
    "Encryptor (antirev-pack.py)  --  Batch encrypt & bundle with topology metadata",
    "",
    "Supporting components:",
    "antirev_client.py  --  Python client: patches import + ctypes.CDLL",
    "build.py  --  Compile/obfuscate Python via Cython, Nuitka, or PyArmor",
    "Diagnostic tools  --  depgraph.py, missing_syms.py",
])

# ── Slide 4: Encryption ─────────────────────────────────────────────
section_slide(prs, 2, "Encryption & Crypto",
              "AES-256-GCM with hardware acceleration")

two_col_slide(prs,
    "AES-256-GCM: Confidentiality + Integrity",
    "Cipher Design",
    [
        "AES-256 in GCM mode (Galois/Counter Mode)",
        "Authenticated encryption: tamper = clean abort",
        "12-byte IV (random per file)",
        "16-byte authentication tag",
        "",
        "Hardware acceleration:",
        "  AES-NI (x86-64) for AES rounds",
        "  PCLMULQDQ for GHASH (carry-less multiply)",
        "  Runtime CPUID detection, software fallback",
        "  ARM64 support via __attribute__((target))",
    ],
    "Two-Pass Strategy",
    [
        "Pass 1 -- GHASH: accumulate auth tag over",
        "  entire ciphertext. Verify BEFORE decrypting.",
        "  Detects bit-flip tampering.",
        "",
        "Pass 2 -- CTR decrypt: counter-mode",
        "  decryption of verified ciphertext.",
        "",
        "Optimization: aes256gcm_onepass()",
        "  combines both passes, halving disk I/O.",
        "",
        "Streaming: 4 MB chunks, never loads",
        "  full ciphertext into memory.",
    ])

# ── Slide 5: Bundle Format ──────────────────────────────────────────
content_slide(prs, "Bundle Format (appended to stub ELF)", [
    "Stub ELF binary  (kernel loads this normally)",
    "+  [num_files: 4B]  [bundle_flags: 1B]",
    "+  Per file:  [name: len-prefixed]  [flags: 1B]  [IV: 12B]  [tag: 16B]  [ct_size: 8B]  [ciphertext]",
    "+  Needed-libs section:  [count: 2B]  [names: len-prefixed]*",
    "+  Trailer:  [bundle_offset: 8B]  [AES key: 32B]  [magic: \"ANTREV01\"]",
    "",
    "Bundle flags control operational mode:",
    "  bit 1 (DAEMON_LIBS)  -- stub is a daemon client, fetch libs from daemon",
    "",
    "Needed-libs section: topologically sorted DT_NEEDED set for the exe",
    "  Computed by antirev-pack.py using Kahn's algorithm",
    "  Stub uses this to create symlinks in the correct order",
])

# ── Slide 6: Stub Launch Flow ───────────────────────────────────────
section_slide(prs, 3, "Stub Launch Flow",
              "From encrypted blob to running process -- no fork, no disk")

diagram_flow(prs, "Stub Execution: Encrypted Binary to Running Process", [
    ("1. Read\nTrailer",
     "pread last 48B\nVerify magic\nExtract key,\nbundle offset"),
    ("2. Scan\nHeaders",
     "Walk bundle\nentries (pread)\nRecord names,\noffsets, sizes"),
    ("3. Decrypt\nto memfd",
     "Stream 4MB\nchunks through\nAES-256-GCM\nWrite to memfd"),
    ("4. Build\nSymlinks",
     "mkdir /tmp/\nantirev_XXX\nln -s /proc/\nself/fd/N lib"),
    ("5. Set\nEnvironment",
     "LD_PRELOAD=\n  exe_shim+\n  dlopen_shim\nLD_LIBRARY_\nPATH+=symdir"),
    ("6. fexecve",
     "Replace stub\nwith decrypted\nexe (same PID)\nKey zeroed"),
], note="Key insight: fexecve replaces the process image in-place. "
        "No fork, no child process, no plaintext on disk. "
        "The kernel loads directly from memfd.")

# ── Slide 7: DT_NEEDED handling ─────────────────────────────────────
two_col_slide(prs,
    "DT_NEEDED Library Handling: Symlink Directory",
    "Problem",
    [
        "Encrypted libs live in memfds (/proc/self/fd/N)",
        "glibc's dynamic linker resolves DT_NEEDED",
        "  by searching directories, not fd numbers",
        "",
        "Old approach: put libs on LD_PRELOAD",
        "  BREAKS symbol resolution order!",
        "  LD_PRELOAD symbols override everything",
        "  Causes silent ABI mismatch bugs",
        "",
        "Need: correct glibc BFS resolution order",
        "  while loading from memfds",
    ],
    "Solution: Symlink Dir",
    [
        "1. Create /tmp/antirev_XXXXXX/",
        "2. For each encrypted DT_NEEDED lib:",
        "     ln -s /proc/self/fd/N  libfoo.so",
        "3. Prepend dir to LD_LIBRARY_PATH",
        "4. glibc's normal BFS finds libs via symlinks",
        "",
        "Result: original symbol lookup order preserved",
        "",
        "LD_PRELOAD carries ONLY the shims:",
        "  exe_shim.so + dlopen_shim.so",
        "  (these intercept /proc/self/exe, dlopen)",
        "",
        "Needed-libs section (topo-sorted) tells stub",
        "  which subset of libs to symlink for each exe",
    ])

# ── Slide 8: exe_shim ───────────────────────────────────────────────
section_slide(prs, 4, "Identity Shims",
              "exe_shim + dlopen_shim: transparent to the application")

content_slide(prs, "exe_shim: Hiding the memfd Identity", [
    "Problem: /proc/self/exe -> memfd:name (deleted) -- application code breaks",
    "",
    "exe_shim intercepts (via LD_PRELOAD):",
    "  readlink(\"/proc/self/exe\")  ->  returns original path on disk",
    "  realpath(\"/proc/self/exe\")  ->  returns original path on disk",
    "  getauxval(AT_EXECFN)          ->  returns original exe filename",
    "  prctl(PR_SET_NAME)              ->  restores process name in ps/top",
    "",
    "Design details:",
    "  Uses raw syscalls (SYS_readlinkat) for fallthrough -- avoids recursion",
    "  Lazy owner detection: is_owner_process() checks /proc/self/exe for memfd",
    "  Child processes inherit LD_PRELOAD but detect non-owner status -> passthrough",
    "",
    "FD cleanup (ANTIREV_CLOSE_FDS):",
    "  After glibc maps DT_NEEDED libs, their memfds are bookkeeping-only",
    "  exe_shim constructor closes them -> frees fd-table slots",
    "  Critical for code using select() with FD_SETSIZE=1024",
])

# ── Slide 9: dlopen_shim ────────────────────────────────────────────
content_slide(prs, "dlopen_shim: Intercepting Dynamic Loading", [
    "Intercepts dlopen(filename, flags) for encrypted libraries",
    "",
    "Eager mode (legacy, small apps):",
    "  Stub pre-fetches all libs, sets ANTIREV_FD_MAP=\"libfoo.so=5,libbar.so=6\"",
    "  dlopen_shim maps basename -> /proc/self/fd/N",
    "",
    "Lazy mode (large apps, daemon):",
    "  Stub keeps daemon socket open (ANTIREV_LIBD_SOCK=fd)",
    "  On dlopen(\"libfoo.so\"):",
    "    1. Send OP_GET_CLOSURE to daemon",
    "    2. Receive lib + all transitive encrypted deps (one round trip)",
    "    3. Materialize symlinks in shared /tmp dir",
    "    4. Preload deps in topological order with RTLD_GLOBAL",
    "    5. real_dlopen(symlink_path) for the root lib",
    "",
    "Why RTLD_GLOBAL: protobuf descriptor dedup across plugins",
    "  Multiple .pb.cc files export same descriptor_table_*",
    "  RTLD_GLOBAL ensures first-loaded copy interposes -> no duplicate abort",
    "",
    "ANTIREV_NO_PRELOAD=1 escape hatch: skip per-dep preload,",
    "  let glibc's natural DT_NEEDED walk run ctors together",
])

# ── Slide 10: Daemon Architecture ───────────────────────────────────
section_slide(prs, 5, "Daemon Architecture",
              "Decrypt once, serve many -- SCM_RIGHTS fd passing")

diagram_flow(prs, "Daemon: Centralized Library Server", [
    ("Startup",
     "Scan dir for\nencrypted .so\nDecrypt all to\nmemfds (threads)"),
    ("Build\nDep Graph",
     "Parse DT_NEEDED\nof each lib\nIntersect with\nencrypted set"),
    ("Bind\nSocket",
     "Abstract Unix\nsocket named by\nkey hash\nH=AES_K(0^128)"),
    ("Serve\nClients",
     "epoll loop\nOP_INIT -> batch\nOP_GET_CLOSURE\n-> topo-sorted"),
    ("Pass FDs",
     "SCM_RIGHTS on\nmsghdr ancillary\nUp to 250 fds\nper OP_BATCH"),
], note="Why daemon? 550 libs x 100 processes = 55,000 memfds without dedup. "
        "Daemon decrypts once, passes fd references. "
        "Socket name derived from key -> zero config, same-key auth.")

content_slide(prs, "Daemon Protocol v2", [
    "Wire format:  [op: 4B]  [payload_len: 4B]  [payload]  [optional SCM_RIGHTS]",
    "",
    "Client -> Daemon:",
    "  OP_INIT (0x01)              request libs (filter or all)",
    "  OP_GET_LIB (0x02)        request single lib",
    "  OP_GET_CLOSURE (0x05)  request lib + transitive encrypted DT_NEEDED closure",
    "  OP_LIST (0x04)              list all lib names (no fds)",
    "  OP_BYE (0x03)              disconnect",
    "",
    "Daemon -> Client:",
    "  OP_BATCH (0x81)    N libs + N fds via SCM_RIGHTS (up to 250 per batch)",
    "  OP_LIB (0x83)         single lib response (status + optional fd)",
    "  OP_NAMES (0x84)    name list (response to OP_LIST)",
    "  OP_END (0x82)         marks end of streaming response",
    "",
    "Closure computation: DFS post-order on dependency graph",
    "  Returns leaves first -> correct load order for dlopen_shim preload loop",
])

# ── Slide 11: Operational Modes ─────────────────────────────────────
section_slide(prs, 6, "Operational Modes",
              "Daemon + Client")

content_slide(prs, "Deployment Model: Daemon + Client", [
    "Lightweight daemon (.antirev-libd)",
    "  Scans its directory for encrypted .so files at startup",
    "  Decrypts 550+ libs once into memfds, serves via abstract socket",
    "  One daemon instance serves 100+ client processes",
    "",
    "Client stubs (protected exes, BFLAG_DAEMON_LIBS)",
    "  Bundle only the exe's ciphertext in the stub trailer",
    "  On launch: connect to daemon, receive lib fds via SCM_RIGHTS",
    "  Lazy fetch: DT_NEEDED set eagerly; dlopen libs fetched on demand",
    "",
    "Massive fd savings compared to per-process lib duplication",
])

# ── Slide 12: Python Integration ────────────────────────────────────
section_slide(prs, 7, "Python Integration",
              "1000+ Python scripts loading encrypted native extensions")

content_slide(prs, "antirev_client.py (Python-level)", [
    "from antirev_client import activate",
    "activate('/path/to/.antirev-libd')",
    "",
    "Patches:",
    "  ctypes.CDLL -> load from memfd",
    "  sys.meta_path -> intercept import",
    "",
    "Pure-Python ELF parser:",
    "  Reads PT_DYNAMIC, DT_NEEDED, DT_SONAME",
    "  No readelf dependency",
    "",
    "Dependency ordering:",
    "  _ensure_loaded() recursively preloads",
    "  transitive DT_NEEDED with RTLD_GLOBAL",
    "",
    "Creates soname symlinks in temp dir prepended to LD_LIBRARY_PATH",
])

# ── Slide 13: Encryptor / Packer ────────────────────────────────────
section_slide(prs, 8, "Encryptor & Packer",
              "antirev-pack.py: batch encrypt with topology metadata")

content_slide(prs, "antirev-pack.py: Configuration-Driven Batch Encryption", [
    "YAML configuration specifies: install_dir, key, blacklist, encrypt_libs",
    "",
    "Phase 1: Discovery",
    "  Walk install_dir, classify ELFs (exe vs lib)",
    "  Apply blacklist (glob/path) and whitelist (encrypt_libs)",
    "",
    "Phase 2: Dependency Analysis",
    "  For each exe, BFS through DT_NEEDED chain",
    "  Classify each dep as encrypted or plaintext",
    "  Topological sort (Kahn's algorithm) among encrypted deps",
    "",
    "Phase 3: Bundle Generation",
    "  Embed needed-libs section (topo-sorted names) in each exe's bundle",
    "  Parallel encryption: ProcessPoolExecutor for CPU-bound AES",
    "  ThreadPoolExecutor for disk I/O (readelf, ldconfig queries)",
    "",
    "Phase 4: Output",
    "  Protected exe = stub + encrypted main + needed-libs section",
    "  Protected daemon = stub + all encrypted libs",
    "  Key stored in 64-hex-char file (chmod 0600)",
])

# ── Slide 14: Key Management ────────────────────────────────────────
content_slide(prs, "Key Management", [
    "Single AES-256 key (32 bytes) per deployment",
    "",
    "Storage:",
    "  Key file: 64 hex characters, chmod 0600",
    "  Embedded in bundle trailer (read by stub at launch)",
    "  Zeroed from memory immediately after use (explicit_bzero)",
    "",
    "Daemon socket naming (zero-config discovery):",
    "  H = AES_K(0^128)   -- encrypt all-zeros block with the key",
    "  Socket = abstract Unix \\0antirev_<first 8 bytes of H as hex>",
    "  Deterministic per key: all stubs with same key find same daemon",
    "  No config files, no port numbers, no environment variables",
    "",
    "Security properties:",
    "  Key never in target process memory post-fexecve",
    "  Plaintext never on disk (memfd only, kernel memory)",
    "  Auth tag verification before any plaintext release",
    "  Same-UID check on daemon socket (SCM_CREDENTIALS)",
])

# ── Slide 15: Diagnostic Tools ──────────────────────────────────────
section_slide(prs, 9, "Diagnostic Tools",
              "Validate before deployment, diagnose in the field")

content_slide(prs, "Diagnostic Tooling", [
    "depgraph.py",
    "  Visualize ELF dependency graph (text or PNG)",
    "  Cycle detection, topological ordering",
    "  --borrowed: find symbols that break under dlopen (implicit deps)",
    "  --find-unresolved: scan dir for missing DT_NEEDED edges",
    "",
    "missing_syms.py",
    "  Scan project dir for all exes/libs with missing DT_NEEDED edges",
    "  Locate providers on LD_LIBRARY_PATH, suggest patchelf --add-needed",
    "  Detect circular dependencies (Tarjan's SCC) and latent cycles",
    "  Combined-graph analysis: catch transitive cycles across multiple fixes",
    "  Per-target duplicate-symbol scan: STRONG x STRONG = error,",
    "    WEAK / mixed = warning; project-scoped by default",
    "  Blacklist third-party dirs, dedup versioned copies",
    "",
    "ANTIREV_DLOPEN_LOG=<path>: runtime dlopen decision logging",
])

# ── Slide 16: Scale & Challenges ────────────────────────────────────
section_slide(prs, 10, "Scale & Challenges",
              "Real-world deployment at enterprise scale")

two_col_slide(prs,
    "Scale & Engineering Challenges",
    "Numbers",
    [
        "100+ executables",
        "550+ shared libraries",
        "1000+ Python scripts",
        "All running concurrently on customer hardware",
        "",
        "FD pressure:",
        "  550 libs x 100 processes = 55,000 memfds",
        "  Daemon reduces to ~550 + overhead",
        "",
        "Startup latency:",
        "  AES-NI: ~1 GB/s throughput",
        "  4 MB streaming chunks (no full load)",
        "  Lazy fetch: only DT_NEEDED set at startup",
        "",
        "select() FD_SETSIZE=1024 limit:",
        "  exe_shim closes DT_NEEDED memfds after load",
    ],
    "Hard Problems Solved",
    [
        "Symbol resolution order: symlink dir preserves",
        "  glibc BFS, unlike LD_PRELOAD approach",
        "",
        "Protobuf descriptor dedup: RTLD_GLOBAL preload",
        "  in closure-deps loop prevents duplicate abort",
        "",
        "Implicit inter-lib deps: ANTIREV_NO_PRELOAD",
        "  escape hatch for natural glibc load semantics",
        "",
        "C++ global ctor ordering: exe_shim lazy detect",
        "  handles calls before shim constructor runs",
        "",
        "Path-based dedup: keep memfd fds open to pin",
        "  /proc/self/fd/N paths, prevent glibc collapse",
        "",
        "32-bit vs 64-bit ldconfig: first-entry-wins",
        "  picks native arch from ldconfig cache",
    ])

# ── Slide 17: Testing ───────────────────────────────────────────────
content_slide(prs, "Testing Strategy", [
    "25+ automated test cases in CMake test suite:",
    "",
    "Functional tests:",
    "  hello (smoke), linked (DT_NEEDED), dlopen (soname), dlopen_nested,",
    "  dlopen_transitive, dlopen_dt_needed, dlopen_reload, dlopen_interpose",
    "",
    "Identity tests:",
    "  proc_self_exe, ctor_readlink, realpath, path_stress",
    "",
    "Process model tests:",
    "  fork_same_lib, fork_diff_lib, script_multi_bin, multi_process",
    "",
    "Daemon/Python tests:",
    "  lib_daemon, python_client_daemon, python_reload",
    "",
    "Security tests:",
    "  wrong_key (clean failure), tamper (bit-flip rejection)",
    "  fd_reduction (FD_SETSIZE compliance)",
    "",
    "Tool tests:",
    "  missing_syms (missing edges, cycles, latent cycles, blacklist, dedup)",
])

# ── Slide 18: Summary ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.7), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

tb = add_textbox(slide, Inches(1.5), Inches(1.2), Inches(7), Inches(1.4))
set_text(tb.text_frame, "Summary", size=40, color=WHITE, bold=True)

tb = add_textbox(slide, Inches(1.5), Inches(3.0), Inches(7.5), Inches(4.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "AntiRev protects business software IP at the deployment layer:",
         size=18, color=LIGHT)
add_para(tf, "", size=8, color=LIGHT)
add_bullet(tf, "AES-256-GCM encryption with hardware acceleration", size=16, color=LIGHT)
add_bullet(tf, "Memory-only execution via memfd + fexecve (no disk plaintext)", size=16, color=LIGHT)
add_bullet(tf, "Daemon architecture for 550+ lib scale (decrypt once, serve many)", size=16, color=LIGHT)
add_bullet(tf, "Symlink-dir approach preserves glibc symbol resolution semantics", size=16, color=LIGHT)
add_bullet(tf, "Transparent to application code (no source changes required)", size=16, color=LIGHT)
add_bullet(tf, "Python integration via antirev_client library", size=16, color=LIGHT)
add_bullet(tf, "Comprehensive diagnostic tooling for pre-deployment validation", size=16, color=LIGHT)

# ── Save ────────────────────────────────────────────────────────────
out = "/mnt/data/anti-rev/antirev_architecture.pptx"
prs.save(out)
print("Saved: %s" % out)
