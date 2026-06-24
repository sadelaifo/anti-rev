#!/usr/bin/env python3
"""Generate the 'kmod2 vs. current (memfd+daemon)' management report deck.

Audience: project manager — focus on security (where the key lives), the
release/CI pipeline, and the client upgrade process. Implementation detail is
kept to what is needed to reason about those.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ── Palette ─────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_CARD = RGBColor(0x25, 0x25, 0x3D)
ACCENT  = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2 = RGBColor(0x90, 0xE0, 0xEF)
ORANGE  = RGBColor(0xFF, 0xA6, 0x2B)
GREEN   = RGBColor(0x06, 0xD6, 0xA0)
RED     = RGBColor(0xEF, 0x47, 0x6F)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xA0, 0xA0, 0xB0)
LIGHT   = RGBColor(0xE0, 0xE0, 0xE8)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return p


def add_para(tf, text, size=16, color=LIGHT, bold=False, space_before=Pt(4),
             space_after=Pt(2), level=0, alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    p.alignment = alignment
    return p


def add_bullet(tf, text, size=15, color=LIGHT, level=0):
    return add_para(tf, text, size=size, color=color, level=level,
                    space_before=Pt(3), space_after=Pt(1))


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD,
                     border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.7), Inches(2), Pt(4))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
    tb = add_textbox(slide, Inches(1.5), Inches(1.0), Inches(7.2), Inches(1.6))
    set_text(tb.text_frame, title, size=38, color=WHITE, bold=True)
    tb = add_textbox(slide, Inches(1.5), Inches(2.9), Inches(7.5), Inches(2.0))
    set_text(tb.text_frame, subtitle, size=18, color=GRAY)
    return slide


def section_slide(prs, number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(1.5), Inches(2))
    set_text(tb.text_frame, "%02d" % number, size=72, color=ACCENT, bold=True)
    tb = add_textbox(slide, Inches(2.5), Inches(1.5), Inches(6.8), Inches(1.2))
    set_text(tb.text_frame, title, size=34, color=WHITE, bold=True)
    if subtitle:
        tb = add_textbox(slide, Inches(2.5), Inches(2.6), Inches(6.8), Inches(1.4))
        set_text(tb.text_frame, subtitle, size=17, color=GRAY)
    return slide


def content_slide(prs, title, bullets, subbullets=None, body_size=15):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD; shape.line.fill.background()
    tb = add_textbox(slide, Inches(0.6), Inches(0.12), Inches(8.9), Inches(0.7))
    set_text(tb.text_frame, title, size=24, color=ACCENT, bold=True)
    tb = add_textbox(slide, Inches(0.6), Inches(1.05), Inches(8.9), Inches(6.2))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for i, b in enumerate(bullets):
        color = LIGHT
        if b.startswith("##"):
            b = b[2:].strip(); color = ORANGE
        if first:
            set_text(tf, b, size=body_size, color=color, bold=(color == ORANGE)); first = False
        else:
            add_bullet(tf, b, size=body_size, color=color, level=0)
            if color == ORANGE:
                tf.paragraphs[-1].font.bold = True
        if subbullets and i in subbullets:
            for sb in subbullets[i]:
                add_bullet(tf, sb, size=body_size - 2, color=GRAY, level=1)
    return slide


def two_col_slide(prs, title, left_title, left_bullets, right_title, right_bullets,
                  left_color=ACCENT, right_color=GREEN, body_size=13):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD; shape.line.fill.background()
    tb = add_textbox(slide, Inches(0.6), Inches(0.12), Inches(8.9), Inches(0.7))
    set_text(tb.text_frame, title, size=24, color=ACCENT, bold=True)

    add_rounded_rect(slide, Inches(0.4), Inches(1.05), Inches(4.5), Inches(6.1), border_color=left_color)
    tb = add_textbox(slide, Inches(0.6), Inches(1.15), Inches(4.2), Inches(0.5))
    set_text(tb.text_frame, left_title, size=16, color=left_color, bold=True)
    tb = add_textbox(slide, Inches(0.6), Inches(1.65), Inches(4.2), Inches(5.4))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(left_bullets):
        (set_text if i == 0 else add_bullet)(tf, b, size=body_size, color=LIGHT)

    add_rounded_rect(slide, Inches(5.1), Inches(1.05), Inches(4.5), Inches(6.1), border_color=right_color)
    tb = add_textbox(slide, Inches(5.3), Inches(1.15), Inches(4.2), Inches(0.5))
    set_text(tb.text_frame, right_title, size=16, color=right_color, bold=True)
    tb = add_textbox(slide, Inches(5.3), Inches(1.65), Inches(4.2), Inches(5.4))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(right_bullets):
        (set_text if i == 0 else add_bullet)(tf, b, size=body_size, color=LIGHT)
    return slide


def _set_cell(cell, text, size=11, color=LIGHT, bold=False, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.07)
    cell.margin_right = Inches(0.07)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def table_slide(prs, title, headers, rows, col_widths=None, font=11,
                header_fill=ACCENT):
    """Comparison table slide. headers = [aspect, colB, colC]; rows = list of
    [aspect, valB, valC]. First column is rendered as an orange label."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD; shape.line.fill.background()
    tb = add_textbox(slide, Inches(0.6), Inches(0.12), Inches(8.9), Inches(0.7))
    set_text(tb.text_frame, title, size=24, color=ACCENT, bold=True)

    nrows = len(rows) + 1
    ncols = len(headers)
    gtable = slide.shapes.add_table(nrows, ncols, Inches(0.3), Inches(1.05),
                                    Inches(9.4), Inches(0.4) * nrows)
    table = gtable.table
    table.first_row = False
    table.horz_banding = False
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        _set_cell(cell, h, size=font + 1, color=BG_DARK, bold=True,
                  align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.LEFT))
    for r, row in enumerate(rows, start=1):
        rowfill = BG_CARD if r % 2 else RGBColor(0x1F, 0x1F, 0x33)
        for j, val in enumerate(row):
            cell = table.cell(r, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = rowfill
            _set_cell(cell, val, size=font,
                      color=(ORANGE if j == 0 else LIGHT), bold=(j == 0))
    return slide


# ── Diagram helpers ─────────────────────────────────────────────────
def diagram_title(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = BG_CARD; bar.line.fill.background()
    tb = add_textbox(slide, Inches(0.6), Inches(0.12), Inches(8.9), Inches(0.7))
    set_text(tb.text_frame, title, size=22, color=ACCENT, bold=True)
    if subtitle:
        tb = add_textbox(slide, Inches(0.5), Inches(6.95), Inches(9.1), Inches(0.5))
        set_text(tb.text_frame, subtitle, size=11, color=GRAY)
    return slide


def dbox(slide, x, y, w, h, title, body="", fill=BG_CARD, border=ACCENT,
         tcolor=None, tsize=13, bsize=10.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border; shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = title; p.font.size = Pt(tsize); p.font.bold = True
    p.font.color.rgb = tcolor or border
    if body:
        for line in body.split("\n"):
            q = tf.add_paragraph(); q.alignment = PP_ALIGN.CENTER
            q.text = line; q.font.size = Pt(bsize); q.font.color.rgb = LIGHT
    return shp


def darrow(slide, x, y, w, h, color=ACCENT2, shape=MSO_SHAPE.DOWN_ARROW):
    a = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background()
    return a


def dlabel(slide, x, y, w, text, color=GRAY, size=10, bold=False):
    tb = add_textbox(slide, Inches(x), Inches(y), Inches(w), Inches(0.3))
    set_text(tb.text_frame, text, size=size, color=color, bold=bold,
             alignment=PP_ALIGN.CENTER)
    return tb


def dconnect(slide, x1, y1, x2, y2, color=ACCENT2, width=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    return c


# ═══════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ── 1. Title ────────────────────────────────────────────────────────
title_slide(prs,
    "Binary Protection: Two Architectures",
    "Current production design  (memfd + daemon + LD_PRELOAD shims)\n"
    "vs.  kmod2  (antirevfs decrypting kernel filesystem)\n\n"
    "Design · CI build output · Runtime footprint · Upgrade process · Key & security posture\n"
    "Prepared for project review — 2026")

# ── 2. Executive summary ────────────────────────────────────────────
content_slide(prs, "Executive Summary (TL;DR)", [
    "##Both designs solve the same problem",
    "Ship 100+ exes / 550+ libs / 1000+ Python scripts to customer-owned hardware so the binaries cannot be reverse-engineered. Plaintext only ever exists in RAM; ciphertext on disk.",
    "",
    "##The two approaches differ in WHERE protection lives",
    "Current: a userspace launcher + a decrypt daemon + injected LD_PRELOAD shims. Each binary is re-wrapped.",
    "kmod2: a kernel filesystem that decrypts files on read. Binaries are untouched; a mount point hides them.",
    "",
    "##The key-position difference is the headline for security",
    "Current: AES key is embedded ONLY inside the launcher/daemon binaries; the 550 lib files on disk are keyless ciphertext — useless if copied.",
    "kmod2: AES key is embedded in EVERY encrypted file. A raw copy of the disk tree is decryptable — protection then rests on the mount + access gate + namespace isolation.",
    "",
    "##Status",
    "Current design = shipping production path. kmod2 = advanced candidate: the full production deployment shape (process-manager stack + access gate + writable views) is now reproduced in automated end-to-end tests on both x86-64 and the ARM64-under-qemu RTOS path — validated on the real x86-64 and SLES targets, but not yet the release path.",
], body_size=14)

# ════════════════ SECTION 1: DESIGN ════════════════
section_slide(prs, 1, "Design",
              "How each architecture protects the binaries")

# Current design
content_slide(prs, "Current design — memfd + daemon + shims", [
    "##Concept: decrypt into anonymous RAM files, run from there",
    "Every protected executable is wrapped in a small C 'stub' launcher. At start it decrypts the real program into a memfd (a kernel in-memory file with no disk path) and runs it via fexecve. No plaintext ever touches the disk.",
    "",
    "##A shared daemon decrypts the 550 libraries once",
    "'lrxd' scans the install tree, decrypts every encrypted .so into memory once, and hands the in-memory file descriptors to each app process over a local socket. 100 processes share one decrypt instead of 55,000 copies.",
    "",
    "##LD_PRELOAD shims keep the app working transparently",
    "Injected helper libraries hide the in-memory identity (so /proc/self/exe still looks normal), redirect dlopen() of encrypted libs, and preserve the normal library-load order.",
    "",
    "##Where the key lives",
    "Embedded in the launcher and daemon binaries only. The library files on disk are keyless ciphertext.",
], body_size=14)

# kmod2 design
content_slide(prs, "kmod2 design — antirevfs kernel filesystem", [
    "##Concept: a filesystem that decrypts on read",
    "Ciphertext lives in a hidden lower tree (.enc/). A small kernel module, antirevfs, is mounted ON TOP of it at the real install path. When any program reads a file, the kernel decrypts it on the fly in the page cache. The app, glibc and Python see plaintext at the real paths — no launcher, no shims, no daemon.",
    "",
    "##Far simpler integration",
    "Business binaries are NOT modified or re-wrapped. Native symbol resolution 'just works', which removes a whole class of dependency-ordering problems the shim design must work around.",
    "",
    "##A kernel access-gate decides who may decrypt",
    "Only authorized processes (an allow-list / signature check) get plaintext through the mount. Copy tools (cp, file managers, backup) are served the ciphertext with the key stripped — so a copy yields useless bytes.",
    "",
    "##Where the key lives",
    "Embedded in every encrypted file. No mount password, no keyring. Convenience cost: the raw .enc/ tree is decryptable if copied (see security section).",
], body_size=13)

# Side by side
COLW = (Inches(1.85), Inches(3.75), Inches(3.8))
table_slide(prs, "Design comparison at a glance",
    ["Aspect", "memfd + daemon  (current)", "kmod2  (antirevfs)"],
    [
        ["Layer", "Userspace only — no kernel changes", "Signed kernel module on the box"],
        ["Binaries", "Each re-wrapped as a stub launcher", "Untouched — mounted, not wrapped"],
        ["Decrypt path", "Daemon decrypts into RAM memfds", "Kernel FS decrypts on read (page cache)"],
        ["Injection", "LD_PRELOAD shims in every process", "None — native loading, real paths"],
        ["/proc view", "memfd:<random> (deleted)", "Real paths — no memfd artifact"],
        ["Dependencies", "Must work around symbol order, dlopen, protobuf dedup", "Native ld.so — problems disappear"],
        ["Key location", "Inside launcher + daemon binaries only", "Embedded in every file"],
        ["Disk if copied", "Keyless ciphertext — safe", "Decryptable — needs gate + ns isolation"],
        ["Status", "Shipping / production", "Candidate; prod shape e2e-tested"],
    ], col_widths=COLW, font=11)

# Key position — security-focused slide for the PM
content_slide(prs, "Where is the AES key? (the security headline)", [
    "##Current design — key is concentrated in executables",
    "The 32-byte AES-256 key is appended inside the launcher of each protected exe and inside the daemon binary. The 550 library files are encrypted WITHOUT the key (keyless container).",
    "A customer copying the lib folder to a USB stick gets keyless ciphertext — undecryptable. The key only travels inside the program binaries themselves.",
    "",
    "##kmod2 — key is embedded in every file",
    "Each encrypted file carries its own key in a trailer; the kernel reads it at decrypt time. There is no separate key file, no password, no keyring to manage.",
    "Trade-off: a raw copy of the .enc/ tree IS decryptable. Protection then depends on (1) the mount stripping the key for unauthorized readers, (2) the access gate, and (3) keeping the .enc/ tree in an isolated mount namespace the operator cannot browse.",
    "",
    "##Both share one hard limit",
    "The CPU runs plaintext, so plaintext is in RAM. Anyone with root on the LIVE box can read it. Our threat model says that's out of scope: the at-box operator is non-technical, and the capable competitor only ever receives copied bits, never the running box.",
    "",
    "##Recommended backstop for both: factory TPM-seal the key per unit",
    "Sealed key never leaves the box and won't unseal on different hardware — exfiltrated ciphertext becomes doubly useless.",
], body_size=13)

# ════════════════ kmod2 DIAGRAMS ════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, BG_DARK)
ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.55), Inches(2), Pt(4))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
tb = add_textbox(slide, Inches(1.2), Inches(2.4), Inches(7.8), Inches(1.2))
set_text(tb.text_frame, "kmod2 in diagrams", size=40, color=WHITE, bold=True)
tb = add_textbox(slide, Inches(1.2), Inches(3.8), Inches(8.2), Inches(1.0))
set_text(tb.text_frame, "Layered architecture · Decrypt read path · Decrypt gate · On-disk container",
         size=18, color=GRAY)

# Fig 1 · layered architecture
slide = diagram_title(prs, "Fig 1 · kmod2 layered architecture (decrypting kernel FS)",
    "Read: passes through antirevfs and is decrypted on the fly; Write: lands in the overlay upper, never touches the lower ciphertext")
dbox(slide, 0.45, 1.15, 5.4, 0.9, "1. App · glibc · ld.so · Python",
     "Sees: real paths + plaintext (identical to unencrypted)", border=GREEN, tcolor=GREEN)
darrow(slide, 3.0, 2.08, 0.3, 0.3); dlabel(slide, 3.4, 2.08, 1.9, "read down: decrypt", ACCENT2)
dbox(slide, 0.45, 2.45, 5.4, 1.45, "2. antirevfs mount (real install path)",
     "read_folio decrypts on read\nkernel page cache shared across processes · real /proc paths\nno memfd · no daemon · no LD_PRELOAD",
     border=ACCENT, tcolor=ACCENT)
darrow(slide, 3.0, 3.95, 0.3, 0.3); dlabel(slide, 3.4, 3.95, 2.0, "write down: overlay", ORANGE)
dbox(slide, 0.45, 4.3, 5.4, 0.8, "3. overlay writable upper",
     "Runtime writes (lock / log / pid) land here, never touch ciphertext", border=ORANGE, tcolor=ORANGE)
darrow(slide, 3.0, 5.15, 0.3, 0.3)
dbox(slide, 0.45, 5.5, 5.4, 1.2, "4. lower .enc/ ciphertext tree (on disk)",
     "ANTREV01 container: magic+iv+tag+ct+key+magic\n(key embedded per file)", border=RED, tcolor=RED)
dbox(slide, 6.1, 1.15, 3.5, 1.7, "Decrypt gate (gate.c)",
     "Authorized process -> plaintext (shared cache)\ncp / backup / file manager -> key-stripped cipher or EACCES\nexec-load by program's own path; data read by caller identity",
     border=ORANGE, tcolor=ORANGE, tsize=12, bsize=10)
dbox(slide, 6.1, 3.0, 3.5, 1.35, "Why it's simpler",
     "Native symbol resolution, no LD_PRELOAD\nno daemon, no memfd fingerprint\ndep-order / protobuf problems vanish",
     border=GREEN, tcolor=GREEN, tsize=12, bsize=10)
dbox(slide, 6.1, 4.5, 3.5, 2.0, "Threat-model tradeoff",
     "Key embedded per file -> raw .enc/ copy is decryptable\nProtection: gate strips key + namespace isolation\nFactory TPM-seal the key as backstop",
     border=RED, tcolor=RED, tsize=12, bsize=10)

# Fig 2 · decrypt read path
slide = diagram_title(prs, "Fig 2 · Decrypt read path: read_folio decrypts on read",
    "First fault decrypts the whole file once (GCM) and verifies the tag, then serves pages from the cache; per-page random decrypt is impossible")
_steps = [
    ("1. read / fault", "app reads a file\nunder the mount\n-> page fault", ACCENT),
    ("2. read_folio", "kernel enters\ndecrypt on read", ACCENT2),
    ("3. get key", "read embedded key\nfrom file trailer\n(zeroed after use)", ORANGE),
    ("4. GCM decrypt", "whole file once\ndecrypt + verify tag", GREEN),
    ("5. page cache", "into per-inode buffer\n-> kernel page cache", ACCENT2),
    ("6. serve plaintext", "serve pages\nshared cross-process", GREEN),
]
_bx, _bw, _gap, _by, _bh = 0.35, 1.42, 0.16, 2.5, 1.95
for _i, (_t, _b, _c) in enumerate(_steps):
    _x = _bx + _i * (_bw + _gap)
    dbox(slide, _x, _by, _bw, _bh, _t, _b, border=_c, tcolor=_c, tsize=12, bsize=10)
    if _i < len(_steps) - 1:
        darrow(slide, _x + _bw + 0.01, _by + _bh / 2 - 0.11, _gap - 0.02, 0.22,
               color=GRAY, shape=MSO_SHAPE.RIGHT_ARROW)
dbox(slide, 0.35, 4.75, 9.27, 1.55, "Unauthorized reader (gate denies)",
     "Skips decrypt -> returns the cipher minus the 40B key trailer (cat sees ANTREV01 magic; objdump can't parse it).\n"
     "GCM authenticates the whole message, so per-page random decrypt is impossible: first fault decrypts the whole file, tag verified once.",
     border=RED, tcolor=RED, tsize=13, bsize=11)

# Fig 3 · decrypt gate decision
slide = diagram_title(prs, "Fig 3 · Decrypt gate: per-process authorization (encrypted files only)",
    "exec-load is authorized by the PROGRAM'S OWN path; everything else (lib load / cp / source) by the caller's current->mm->exe_file")
dbox(slide, 3.1, 1.05, 3.8, 0.55, "open an encrypted file (ANTREV01) under the mount", border=ACCENT, tsize=12)
darrow(slide, 4.85, 1.62, 0.3, 0.26)
dbox(slide, 4.0, 1.92, 2.0, 0.92, "exec-load ?", border=ACCENT2, tcolor=WHITE,
     tsize=13, shape=MSO_SHAPE.DIAMOND)
dconnect(slide, 4.0, 2.38, 2.05, 3.1); dlabel(slide, 2.3, 2.42, 1.6, "yes · exec", ACCENT2, 10, True)
dbox(slide, 1.0, 3.1, 2.1, 0.95, "program itself\nwhitelisted ?", border=ACCENT2, tcolor=WHITE,
     tsize=12, shape=MSO_SHAPE.DIAMOND)
dconnect(slide, 1.55, 4.05, 1.15, 4.45); dlabel(slide, 0.2, 4.1, 1.0, "yes", GREEN, 10, True)
dbox(slide, 0.35, 4.45, 1.6, 0.6, "OK: exec", border=GREEN, tcolor=GREEN, tsize=12)
dconnect(slide, 2.55, 4.05, 3.0, 4.45); dlabel(slide, 3.0, 4.1, 1.0, "no", RED, 10, True)
dbox(slide, 2.15, 4.45, 1.75, 0.6, "EACCES (126)", border=RED, tcolor=RED, tsize=12)
dconnect(slide, 6.0, 2.38, 7.2, 3.1); dlabel(slide, 6.05, 2.42, 2.0, "no · data read", ACCENT2, 10, True)
dbox(slide, 6.0, 3.1, 2.4, 0.95, "caller process\nwhitelisted ?", border=ACCENT2, tcolor=WHITE,
     tsize=12, shape=MSO_SHAPE.DIAMOND)
dconnect(slide, 6.6, 4.05, 5.7, 4.45); dlabel(slide, 4.95, 4.1, 1.0, "yes", GREEN, 10, True)
dbox(slide, 4.6, 4.45, 2.05, 0.78, "decrypt -> plaintext\n(shared page cache)", border=GREEN, tcolor=GREEN,
     tsize=12, bsize=10)
dconnect(slide, 7.9, 4.05, 8.5, 4.45); dlabel(slide, 8.5, 4.1, 1.0, "no", RED, 10, True)
dbox(slide, 7.6, 4.45, 2.0, 0.9, "gate_passthrough\n_cipher ?", border=ACCENT2, tcolor=WHITE,
     tsize=11, shape=MSO_SHAPE.DIAMOND)
dconnect(slide, 8.2, 5.35, 7.75, 5.7); dlabel(slide, 6.75, 5.42, 1.0, "yes", ORANGE, 10, True)
dbox(slide, 6.9, 5.7, 1.7, 0.55, "key-stripped cipher", border=ORANGE, tcolor=ORANGE, tsize=10)
dconnect(slide, 9.0, 5.35, 9.2, 5.7); dlabel(slide, 8.7, 5.42, 0.8, "no", RED, 10, True)
dbox(slide, 8.65, 5.7, 1.1, 0.55, "EACCES", border=RED, tcolor=RED, tsize=11)

# Fig 4 · on-disk container format
slide = diagram_title(prs, "Fig 4 · On-disk container: embedded-key trailer (no mount key)",
    "Authorized reader: kernel reads the trailer key and decrypts on the fly (stack copy zeroed after use); unauthorized reader: the 40B trailer is stripped")
_segs = [("ANTREV01\n8B", 1.1, ACCENT), ("IV\n12B", 1.0, ACCENT2),
         ("TAG\n16B", 1.2, ORANGE), ("ciphertext ct\n(plaintext, GCM)", 3.0, GREEN),
         ("KEY\n32B", 1.5, RED), ("ANTREV01\n8B", 1.1, ACCENT)]
_sx, _sy, _sh = 0.55, 2.7, 1.05
_xs = []; _cx = _sx
for (_t, _w, _c) in _segs:
    _xs.append((_cx, _w, _c, _t)); _cx += _w
for (_x, _w, _c, _t) in _xs:
    dbox(slide, _x, _sy, _w, _sh, _t, border=_c, tcolor=_c, tsize=11)
dlabel(slide, _xs[0][0], 2.28, (_xs[2][0] + _xs[2][1]) - _xs[0][0],
       "HDR = 36B (magic+iv+tag)", ACCENT2, 11, True)
dlabel(slide, _xs[3][0], 2.28, _xs[3][1], "plaintext_len = file - 76B", GREEN, 11, True)
dlabel(slide, _xs[4][0], 2.28, (_xs[5][0] + _xs[5][1]) - _xs[4][0],
       "TRAILER = 40B (key+magic)", RED, 11, True)
dlabel(slide, _xs[4][0], _sy + _sh + 0.05, (_xs[5][0] + _xs[5][1]) - _xs[4][0],
       "unauthorized reader -> drop this 40B -> keyless container", RED, 10.5, True)
dbox(slide, 0.55, 4.55, 9.05, 1.75, "Key points",
     "Key travels with the file: a raw .enc/ copy is decryptable -> requires namespace isolation + factory TPM-seal as backstop.\n"
     "Authorized reader: reads the trailer key and decrypts on the fly.\n"
     "Unauthorized reader: the 40B trailer is stripped; cp / vim / objdump only get a useless keyless container (objdump: format not recognized).",
     border=ACCENT, tcolor=ACCENT, tsize=13, bsize=11)

# ════════════════ SECTION 2: CI BUILD OUTPUT ════════════════
section_slide(prs, 2, "After the CI build",
              "What the software package contains — new files & modified files")

content_slide(prs, "Release pipeline — two stages (both designs)", [
    "##Stage A — Compile (CI, key-independent, reusable)",
    "Builds the protection toolchain from source. Produces architecture binaries with NO customer key in them. The same Stage-A output is reused across all deployments.",
    "  Current: stub launcher, antirev_shim_<arch>.so, lrxd-<arch> daemon, packer tools",
    "  kmod2: antirevfs.ko (one per kernel version, via DKMS), mount tools",
    "",
    "##Stage B — Pack (per-deployment, uses the secret key)",
    "Takes the business install tree + the deployment key + Stage-A artifacts and produces the encrypted, shippable package. This is the step that injects/uses the key and yields the client deliverable.",
    "  Current: antirev-pack.py  → re-wrapped exes + encrypted libs + daemon",
    "  kmod2: antirev-fs-pack.py → ciphertext .enc/ lower tree + manifest",
    "",
    "##Why this split matters to release management",
    "Stage A is ordinary CI (build once, sign, archive). Stage B is the controlled, key-bearing step that should run in a protected release environment. The key is the only secret; everything else is reproducible from source.",
], body_size=14)

table_slide(prs, "After CI: what the package looks like",
    ["Aspect", "memfd + daemon  (current)", "kmod2  (antirevfs)"],
    [
        ["New binaries", "lrxd-<arch> daemon, antirev_shim_<arch>.so, LD_PRELOAD launch scripts", "antirevfs.ko (signed), mount tools, boot/systemd mount unit"],
        ["New data files", "—", ".enc/ ciphertext tree, manifest.json, proj-pack.yaml, /etc/authorized_apps.txt"],
        ["Modified", "Every exe → stub+ciphertext+key; every .so → keyless ciphertext", "None — install path becomes a MOUNT POINT over .enc/"],
        ["Python loaders", "Use antirev_client to load encrypted .so", "Unchanged — load natively through the mount"],
        ["Key artifact", "Separate 64-hex key file (0600) / TPM-sealed", "Embedded per file — no separate key to ship"],
        ["3rd-party libs", "Untouched (coexist plaintext)", "Left plaintext in the tree (blacklisted)"],
    ], col_widths=COLW, font=11)

# ════════════════ SECTION 3: RUNTIME ════════════════
section_slide(prs, 3, "At runtime",
              "What the install directory and the running processes look like")

table_slide(prs, "Runtime footprint on the box",
    ["Aspect", "memfd + daemon  (current)", "kmod2  (antirevfs)"],
    [
        ["On disk", "Stub-wrapped exes + keyless ciphertext libs + lrxd daemon", ".enc/ lower tree (hidden) + mount point showing real paths"],
        ["Extra process", "One long-lived lrxd daemon", "None"],
        ["Injected libs", "LD_PRELOAD shims in every process", "None — apps run unmodified"],
        ["/proc/<pid>/exe", "memfd:<random> (deleted)", "Real path"],
        ["/proc maps", "memfd entries, no library paths", "Real library paths"],
        ["Other artifacts", "Abstract Unix socket, /tmp/antirev_* symlink dirs, open memfds", "lsmod antirevfs, mount entries, shared page cache, kernel gate"],
        ["Writable bin/lib", "Normal files", "Overlay upper stacked over read-only antirevfs"],
    ], col_widths=COLW, font=11)

content_slide(prs, "Runtime: the writable-view detail (kmod2)", [
    "##The mount is read-only by design",
    "antirevfs only serves shipped, encrypted content. But real business software writes lock / pid / log files right next to its binaries (e.g. the GUI writes QtApplication.pid into bin/).",
    "",
    "##Fix: a stock overlay provides a writable top layer",
    "antirev-mount-rw stacks a normal writable layer over the decrypting filesystem. Decrypted library reads fall through to antirevfs; the app's runtime writes land in the writable upper and NEVER touch the ciphertext .enc/ tree. Both bin/ and lib/ are mounted this way in the reference deployment.",
    "",
    "##Net operator-visible picture",
    "A normal-looking bin/ and lib/ that the app reads and writes as usual — while the real bytes on disk stay encrypted, and only authorized processes can pull plaintext through the mount.",
], body_size=14)

# ── Developer self-verification scenario ────────────────────────────
table_slide(prs, "Developer self-test: impact of building a new binary",
    ["Aspect", "memfd + daemon  (current)", "kmod2  (antirevfs)"],
    [
        ["To test a new binary", "Pack first (encrypt + wrap stub) → start daemon → run", "Encrypt into .enc/ → mount → run (or develop plaintext, encrypt once to verify)"],
        ["Same as compiled?", "Re-wrapped — not byte-identical; process model changes", "Byte-identical once decrypted — fully transparent"],
        ["Debugging (gdb/core)", "Hard — /proc shows memfd, readlink spoofed, symbols hidden, core points at memfd", "Easy — real paths, native symbols, gdb/core dumps work"],
        ["Protection-only bugs", "A class appears only when protected (symbol order, protobuf dedup, implicit deps / NO_PRELOAD)", "Essentially none — native loading, plaintext == encrypted behavior"],
        ["Privileges", "Userspace — no root needed", "Root (insmod + mount); dev can set gate_enforce=0"],
        ["Iteration loop", "Edit → re-pack → restart daemon → run", "Edit → re-encrypt that file → remount → run"],
        ["Adding a dependency", "Update needed-libs metadata / re-pack the exe", "None — native DT_NEEDED resolution"],
        ["Gate impact on dev", "No gate concept", "If gate on, allow-list shell/interpreter/gdb by basename"],
        ["Self-test in CI", "Pure userspace — runs on a plain runner", "Needs privileged container / kernel module — heavier CI"],
        ["Encryption visible to dev?", "High — wrapping / daemon / memfd artifacts all show", "Low — real paths; looks like an ordinary directory"],
    ], col_widths=(Inches(1.95), Inches(3.7), Inches(3.75)), font=10)

content_slide(prs, "Developer self-test: verifying without revealing encryption", [
    "##The need: a dev edits their own exe/.so, drops it on the FS to verify — and other developers must NOT notice the package is encrypted",
    "",
    "##Shared baseline (both designs): devs develop & self-test on the plaintext tree",
    "Encryption happens only at the CI pack step (Stage B). Developers build, run and debug on the unencrypted source/build tree — identical to an unprotected system, so encryption is invisible. This is the default and answers most self-test needs.",
    "",
    "##When the change must run INSIDE the protected package, the two diverge sharply:",
    "kmod2: real paths, transparent — drop the self-built binary into the writable overlay (lands as plaintext) and it coexists with the decrypted libs at real paths. The dev sees an ordinary directory: no memfd, no daemon, no LD_PRELOAD — encryption is imperceptible.",
    "  Caveat: with the gate on, a non-whitelisted self-built binary reading encrypted libs gets EACCES (looks like a plain permission error — which can itself raise suspicion) → use gate_enforce=0 on dev boxes, or whitelist its basename temporarily.",
    "memfd + daemon: to use the encrypted libs the self-built binary must launch via the protected start script (LD_PRELOAD + daemon socket); and the stub wrapping, lrxd daemon and memfd artifacts are all exposed → hard to hide encryption from the dev. In practice devs verify on the plaintext tree; the protected form is verified by CI / a dedicated team.",
    "",
    "##Bottom line",
    "'Keep encryption imperceptible to developers' is met natively by kmod2 at the FS layer (transparent, real paths); memfd + daemon, whose machinery is externally visible, can essentially only achieve it by 'devs only ever touch the plaintext tree'.",
], body_size=13)

content_slide(prs, "How each keeps encryption imperceptible: two recipes", [
    "##kmod2 recipe (protected state is itself transparent — just configure it)",
    "Mount the package as antirevfs + overlay-rw + passdata at the real install path; gate_enforce=0 on dev boxes.",
    "Self-replace: cp the self-built binary into the mount → lands in the overlay upper (plaintext, served verbatim, ungated); a same-named file shadows the decrypted lower lib, so the dev's version takes effect immediately.",
    "Self-verify: launch normally — real paths, no memfd / no daemon / no LD_PRELOAD, behavior == plaintext.",
    "",
    "##memfd + daemon recipe (machinery is visible — keep the dev on the plaintext plane)",
    "Main path: devs replace + verify only on the unencrypted build tree — no encryption in their world, so imperceptible; encryption happens only at CI Stage B, the protected form verified by CI / a dedicated team.",
    "Cost: misses 'works plaintext, breaks encrypted' bugs (symbol order / protobuf / implicit deps / NO_PRELOAD) — CI must catch them.",
    "If it must run protected: a wrapper script hides the LD_PRELOAD / daemon-socket env → hides the commands, but /proc memfd and ps lrxd stay visible; and a plaintext same-name replacement is shadowed by the daemon's symlink dir.",
], body_size=13)

table_slide(prs, "Imperceptible encryption: the two designs",
    ["Aspect", "memfd + daemon  (current)", "kmod2  (antirevfs)"],
    [
        ["Where it hides", "Can't hide runtime — keep the dev on the plaintext tree", "In the kernel read path — protected state is itself transparent"],
        ["Self-replace", "cp on plaintext tree; a same-name encrypted lib is shadowed by the daemon", "cp into mount → overlay upper plaintext; same name auto-shadows, instant"],
        ["Self-verify scope", "Functional only; protection-only bugs need CI / a team", "Verify directly on real paths; behavior == plaintext"],
        ["Required step", "Wrapper script to hide env (still can't hide memfd / daemon)", "gate_enforce=0 on dev boxes (or basename whitelist)"],
        ["Residual exposure", "/proc memfd, ps lrxd, LD_PRELOAD all show", "Only lsmod / mount reveal the module; normal workflow imperceptible"],
    ], col_widths=COLW, font=11)

# ════════════════ SECTION 4: PARTIAL UPGRADE ════════════════
section_slide(prs, 4, "Partial upgrade",
              "Replacing / adding only some libraries in the field")

table_slide(prs, "Partial upgrade — replacing / adding some libraries",
    ["Step", "memfd + daemon  (current)", "kmod2  (antirevfs)"],
    [
        ["Key requirement", "MUST reuse the one deployment key (all files + daemon share it)", "Files self-describe their key; reuse project key for consistency"],
        ["Re-encrypt", "Changed/new .so with the existing key", "Changed/new ELF into .enc/ (embedded-key form)"],
        ["Place files", "Drop ciphertext libs into the tree", "Swap the file(s) in the .enc/ lower tree"],
        ["New dep edge", "Re-pack the affected exe (needed-libs list changes)", "Not needed — native resolution"],
        ["Activate", "Restart the lrxd daemon (re-scan & re-decrypt)", "REMOUNT the tree (classification cached per file)"],
        ["Gotcha", "Daemon must be restarted", "Needs a quiet stack — mapped libs pin the module"],
    ], col_widths=COLW, font=11)

content_slide(prs, "Partial upgrade — key takeaways for release planning", [
    "##The key is the one artifact that MUST survive every upgrade",
    "Current design: mandatory — every file in the deployment is bound to one key; new libs must be encrypted with it or the daemon can't serve a consistent set.",
    "kmod2: less strict — files self-describe their key — but keeping one project keyfile keeps the manifest and tooling clean.",
    "",
    "##Practical guidance",
    "Store the per-deployment key in the release vault AND TPM-seal a copy on each unit at the factory. Partial patches then never need to ship a key — the technician re-encrypts against the retained key and drops files in.",
    "",
    "##The one operational gotcha per design",
    "Current: remember to RESTART the daemon (it decrypts at startup).",
    "kmod2: remember to REMOUNT (inode classification is cached); a quiet stack is needed because mapped libraries pin the module.",
], body_size=14)

# ════════════════ SECTION 5: FULL UPGRADE ════════════════
section_slide(prs, 5, "Full upgrade",
              "Replacing the entire protected software set")

table_slide(prs, "Full upgrade — replacing the entire protected set",
    ["Step", "memfd + daemon  (current)", "kmod2  (antirevfs)"],
    [
        ["Stop", "Business stack + lrxd daemon", "Business stack (mapped .so pin the module)"],
        ["Tear down", "—", "Bring mounts down (antirev-mount-rw --down)"],
        ["Key choice", "Keep (drop-in) or rotate (re-pack all together, re-seal TPM)", "Keep or rotate; each file carries its own key"],
        ["Deploy", "Swap packed tree (exes, libs, daemon, shims, scripts)", "Replace .enc/ lower tree + manifest"],
        ["Kernel dep.", "None — pure userspace swap", ".ko must match kernel; DKMS rebuild + Secure Boot signing"],
        ["Activate", "Start daemon, then the stack", "Remount (antirev-remount-proj.sh), start the stack"],
        ["Rollback", "Keep previous packed tree + key; swap back", "Keep previous .enc/ tree; remount"],
    ], col_widths=COLW, font=11)

content_slide(prs, "Full upgrade — what release management must own", [
    "##Current design — fully self-contained, no host coupling",
    "A full upgrade is a userspace file swap plus a daemon restart. Nothing depends on the customer's kernel. Simplest to support across heterogeneous customer hardware.",
    "",
    "##kmod2 — adds a kernel-module lifecycle to own",
    "The .ko must be built and signed for each customer kernel version. DKMS rebuilds it automatically on kernel upgrade, but Secure Boot requires our module-signing key to be enrolled. This is new release/QA surface: a kernel update on the customer box can require a module rebuild.",
    "",
    "##Recommendation for either path",
    "Ship full upgrades as a versioned, signed bundle (Stage-A artifacts + Stage-B packed tree) with a tested rollback (retain the previous tree + key). Treat the deployment key as a release secret with factory TPM-sealing as the field backstop.",
], body_size=14)

# ════════════════ SECTION 6: SECURITY / WRAP-UP ════════════════
section_slide(prs, 6, "Security posture, trade-offs & recommendation",
              "Everything the project decision hinges on")

content_slide(prs, "Threat model (agreed) — keep it in view", [
    "##At-box adversary = the client/operator",
    "Physically owns the hardware but is NON-TECHNICAL: GUI buttons, drag-and-drop, copy-to-USB. No concept of root, mounts, decryption, or memory forensics. Attack = 'copy the software folder and hand it to someone.'",
    "",
    "##Capable adversary = the competitor",
    "Expert reverse-engineers, but DECOUPLED from the live appliance — they only ever receive the bits the operator copied, never root on a running box.",
    "",
    "##Therefore the whole problem reduces to:",
    "'What a non-technical user can copy via GUI/USB must be ciphertext.'",
    "Sophisticated root-on-live-box attacks are genuinely out of scope: the person at the box can't do them, and the person who can isn't at the box.",
], body_size=15)

table_slide(prs, "Security scorecard against the threat model",
    ["Aspect", "memfd + daemon  (current)", "kmod2  (antirevfs)"],
    [
        ["Copy lib folder", "Keyless ciphertext — key is NOT in the libs (strong)", "Via mount: key-stripped ciphertext (useless)"],
        ["Copy raw disk tree", "Keyless ciphertext — safe", ".enc/ is DECRYPTABLE — must hide in mount namespace"],
        ["Plaintext exposure", "RAM memfds; memfd artifact hints 'protected'", "Kernel page cache; real paths, no artifact"],
        ["Residual risk", "Key concentrated in exe + daemon binaries", "Embedded per-file key is the weak point"],
        ["Hardening path", "Obfuscation + TPM seal", "Wrap / TPM-seal the per-file key (TODO)"],
        ["Maturity", "Production-proven at full scale", "Candidate; prod stack shape now e2e-tested"],
    ], col_widths=COLW, font=11)

content_slide(prs, "Recommendation & next steps", [
    "##Keep the current memfd+daemon design as the shipping path",
    "It is production-proven, fully userspace (no customer-kernel coupling), and its at-rest posture is strong: the bulk of the deployment (550 libs) is keyless ciphertext.",
    "",
    "##Continue kmod2 as the strategic candidate",
    "It is dramatically simpler to integrate (no shims, no dependency-ordering workarounds, binaries untouched) and leaves no memfd fingerprint. The real production deployment shape — process-manager stack, access gate, writable views — is now reproduced in automated end-to-end tests (x86-64 and ARM64-under-qemu). Its main gaps are the embedded key + the kernel-module lifecycle.",
    "",
    "##To make kmod2 release-ready, close these:",
    "  1. Harden the embedded key — obfuscate / TPM-seal the per-file trailer so a copied .enc/ tree is also undecryptable",
    "  2. Module signing for Secure Boot + DKMS .deb/.rpm packaging",
    "  3. Promote the access-gate from allow-list (demo) to signature-based authorization",
    "  4. Validate the Docker / container-namespace path for the ARM64-on-x86 RTOS slave",
    "",
    "##For BOTH: factory TPM-seal the per-unit key",
    "This is the single highest-leverage control — it makes any exfiltrated ciphertext useless off the original hardware.",
], body_size=13)

# Closing
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.7), Inches(2), Pt(4))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
tb = add_textbox(slide, Inches(1.5), Inches(1.2), Inches(7.5), Inches(1.4))
set_text(tb.text_frame, "One-line summary", size=34, color=WHITE, bold=True)
tb = add_textbox(slide, Inches(1.5), Inches(3.0), Inches(7.8), Inches(4.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Current design = proven, userspace, key hidden in executables, libs keyless — keep shipping it.",
         size=18, color=LIGHT)
add_para(tf, "", size=8)
add_para(tf, "kmod2 = simpler & cleaner, key embedded per file — promising once the key is hardened and the kernel module is signed/packaged.",
         size=18, color=LIGHT)
add_para(tf, "", size=8)
add_para(tf, "For both: the deployment key is the crown jewel — vault it in release, TPM-seal it per unit in the factory.",
         size=18, color=ACCENT2, bold=True)

# ── Save ────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "kmod2_vs_daemon_report.pptx")
prs.save(out)
print("Saved:", out, "—", len(prs.slides._sldIdLst), "slides")
