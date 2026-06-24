#!/usr/bin/env python3
"""生成「kmod2 对比 现行方案（memfd + 守护进程）」管理层汇报 PPT（简体中文）。

面向项目经理 —— 聚焦安全（密钥位置）、发布/CI 流程、客户升级流程。
实现细节仅保留推理这些问题所必需的部分。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ── 配色 ─────────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_CARD = RGBColor(0x25, 0x25, 0x3D)
ACCENT  = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2 = RGBColor(0x90, 0xE0, 0xEF)
ORANGE  = RGBColor(0xFF, 0xA6, 0x2B)
GREEN   = RGBColor(0x06, 0xD6, 0xA0)
RED     = RGBColor(0xEF, 0x47, 0x6F)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xA0, 0xA0, 0xB0)
LIGHT   = RGBColor(0xE0, 0xE0, 0xE8)

CJK = "Microsoft YaHei"   # 简体中文字体


def _apply_cjk(font):
    """让一个 font 对象同时为拉丁/东亚/复杂文种设置中文字体。"""
    font.name = CJK
    rPr = font._rPr
    if rPr is None:
        return
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", CJK)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    _apply_cjk(p.font)
    return p


def add_para(tf, text, size=16, color=LIGHT, bold=False, space_before=Pt(4),
             space_after=Pt(2), level=0, alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    p.alignment = alignment
    _apply_cjk(p.font)
    return p


def add_bullet(tf, text, size=15, color=LIGHT, level=0):
    return add_para(tf, text, size=size, color=color, level=level,
                    space_before=Pt(3), space_after=Pt(1))


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD,
                     border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.7), Inches(2), Pt(4))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
    tb = add_textbox(slide, Inches(1.5), Inches(1.0), Inches(7.2), Inches(1.6))
    set_text(tb.text_frame, title, size=36, color=WHITE, bold=True)
    tb = add_textbox(slide, Inches(1.5), Inches(2.9), Inches(7.6), Inches(2.4))
    set_text(tb.text_frame, subtitle, size=17, color=GRAY)
    return slide


def section_slide(prs, number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(1.5), Inches(2))
    set_text(tb.text_frame, "%02d" % number, size=72, color=ACCENT, bold=True)
    tb = add_textbox(slide, Inches(2.5), Inches(1.5), Inches(6.8), Inches(1.2))
    set_text(tb.text_frame, title, size=34, color=WHITE, bold=True)
    if subtitle:
        tb = add_textbox(slide, Inches(2.5), Inches(2.6), Inches(6.8), Inches(1.4))
        set_text(tb.text_frame, subtitle, size=17, color=GRAY)
    return slide


def content_slide(prs, title, bullets, subbullets=None, body_size=15):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD; shape.line.fill.background()
    tb = add_textbox(slide, Inches(0.6), Inches(0.12), Inches(8.9), Inches(0.7))
    set_text(tb.text_frame, title, size=24, color=ACCENT, bold=True)
    tb = add_textbox(slide, Inches(0.6), Inches(1.05), Inches(8.9), Inches(6.2))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for i, b in enumerate(bullets):
        color = LIGHT
        if b.startswith("##"):
            b = b[2:].strip(); color = ORANGE
        if first:
            set_text(tf, b, size=body_size, color=color, bold=(color == ORANGE)); first = False
        else:
            add_bullet(tf, b, size=body_size, color=color, level=0)
            if color == ORANGE:
                tf.paragraphs[-1].font.bold = True
        if subbullets and i in subbullets:
            for sb in subbullets[i]:
                add_bullet(tf, sb, size=body_size - 2, color=GRAY, level=1)
    return slide


def _set_cell(cell, text, size=11, color=LIGHT, bold=False, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.07)
    cell.margin_right = Inches(0.07)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        _apply_cjk(run.font)


def table_slide(prs, title, headers, rows, col_widths=None, font=11,
                header_fill=ACCENT):
    """对比表幻灯片。headers = [对比项, 列B, 列C]；rows = [对比项, 值B, 值C] 列表。
    第一列以橙色标签呈现。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD; shape.line.fill.background()
    tb = add_textbox(slide, Inches(0.6), Inches(0.12), Inches(8.9), Inches(0.7))
    set_text(tb.text_frame, title, size=24, color=ACCENT, bold=True)

    nrows = len(rows) + 1
    ncols = len(headers)
    gtable = slide.shapes.add_table(nrows, ncols, Inches(0.3), Inches(1.05),
                                    Inches(9.4), Inches(0.4) * nrows)
    table = gtable.table
    table.first_row = False
    table.horz_banding = False
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        _set_cell(cell, h, size=font + 1, color=BG_DARK, bold=True)
    for r, row in enumerate(rows, start=1):
        rowfill = BG_CARD if r % 2 else RGBColor(0x1F, 0x1F, 0x33)
        for j, val in enumerate(row):
            cell = table.cell(r, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = rowfill
            _set_cell(cell, val, size=font,
                      color=(ORANGE if j == 0 else LIGHT), bold=(j == 0))
    return slide


# ── 图解辅助 ─────────────────────────────────────────────────────────
def diagram_title(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = BG_CARD; bar.line.fill.background()
    tb = add_textbox(slide, Inches(0.6), Inches(0.12), Inches(8.9), Inches(0.7))
    set_text(tb.text_frame, title, size=22, color=ACCENT, bold=True)
    if subtitle:
        tb = add_textbox(slide, Inches(0.5), Inches(6.95), Inches(9.1), Inches(0.5))
        set_text(tb.text_frame, subtitle, size=11, color=GRAY)
    return slide


def dbox(slide, x, y, w, h, title, body="", fill=BG_CARD, border=ACCENT,
         tcolor=None, tsize=13, bsize=10.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border; shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = title; p.font.size = Pt(tsize); p.font.bold = True
    p.font.color.rgb = tcolor or border; _apply_cjk(p.font)
    if body:
        for line in body.split("\n"):
            q = tf.add_paragraph(); q.alignment = PP_ALIGN.CENTER
            q.text = line; q.font.size = Pt(bsize); q.font.color.rgb = LIGHT
            _apply_cjk(q.font)
    return shp


def darrow(slide, x, y, w, h, color=ACCENT2, shape=MSO_SHAPE.DOWN_ARROW):
    a = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background()
    return a


def dlabel(slide, x, y, w, text, color=GRAY, size=10, bold=False):
    tb = add_textbox(slide, Inches(x), Inches(y), Inches(w), Inches(0.3))
    set_text(tb.text_frame, text, size=size, color=color, bold=bold,
             alignment=PP_ALIGN.CENTER)
    return tb


def dconnect(slide, x1, y1, x2, y2, color=ACCENT2, width=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    return c


# ═══════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

COLW = (Inches(1.85), Inches(3.75), Inches(3.8))
HDR = ["对比项", "memfd + 守护进程（现行）", "kmod2（antirevfs）"]

# ── 1. 标题 ──────────────────────────────────────────────────────────
title_slide(prs,
    "二进制保护：两种架构方案",
    "现行生产方案（memfd + 守护进程 + LD_PRELOAD 注入）\n"
    "对比  kmod2（antirevfs 解密内核文件系统）\n\n"
    "设计 · CI 构建产物 · 运行时形态 · 升级流程 · 密钥与安全态势\n"
    "面向项目评审 —— 2026")

# ── 2. 执行摘要 ─────────────────────────────────────────────────────
content_slide(prs, "执行摘要（要点速览）", [
    "##两种方案解决的是同一个问题",
    "把 100+ 可执行文件 / 550+ 动态库 / 1000+ Python 脚本交付到客户自有硬件上，使二进制无法被逆向。明文只存在于内存中，磁盘上只有密文。",
    "",
    "##差异在于「保护逻辑位于何处」",
    "现行方案：用户态启动器 + 解密守护进程 + 注入的 LD_PRELOAD shim；每个二进制都被重新封装。",
    "kmod2：一个在读取时解密文件的内核文件系统；二进制原封不动，仅用挂载点遮蔽。",
    "",
    "##密钥位置是安全方面的核心结论",
    "现行方案：AES 密钥仅嵌入在启动器/守护进程二进制中；磁盘上 550 个库文件是无密钥密文 —— 即使被拷走也无法解密。",
    "kmod2：AES 密钥嵌入在每个加密文件中。直接拷贝磁盘目录即可解密 —— 保护因此依赖挂载层 + 访问门禁 + 命名空间隔离。",
    "",
    "##现状",
    "现行方案 = 正在交付的生产路径。kmod2 = 进阶候选方案：完整生产部署形态（进程管理器栈 + 访问门禁 + 可写视图）已在 x86-64 与 ARM64-under-qemu 路径上有自动化端到端测试覆盖；已在真实 x86-64 与 SLES 目标上验证，但尚未作为发布路径。",
], body_size=14)

# ════════════════ 第 1 部分：设计 ════════════════
section_slide(prs, 1, "设计", "两种架构各自如何保护二进制")

content_slide(prs, "现行方案 —— memfd + 守护进程 + shim", [
    "##理念：解密到匿名内存文件中并从那里运行",
    "每个受保护的可执行文件都被一个小型 C「stub」启动器封装。启动时它把真正的程序解密到 memfd（一个没有磁盘路径的内核内存文件）中，并通过 fexecve 运行。磁盘上永远不会出现明文。",
    "",
    "##一个共享守护进程一次性解密 550 个库",
    "「lrxd」扫描安装目录，把每个加密的 .so 在内存中解密一次，再通过本地 socket 把内存文件描述符分发给每个应用进程。100 个进程共享一次解密，而非 55,000 份拷贝。",
    "",
    "##LD_PRELOAD shim 让应用透明运行",
    "注入的辅助库隐藏内存身份（使 /proc/self/exe 看起来正常）、重定向加密库的 dlopen()、并保持正常的库加载顺序。",
    "",
    "##密钥位置",
    "仅嵌入在启动器与守护进程二进制中。磁盘上的库文件是无密钥密文。",
], body_size=14)

content_slide(prs, "kmod2 方案 —— antirevfs 内核文件系统", [
    "##理念：一个在读取时解密的文件系统",
    "密文存放在隐藏的下层目录（.enc/）中。一个小型内核模块 antirevfs 挂载在其之上、位于真实安装路径。任何程序读取文件时，内核在页缓存中即时解密。应用、glibc 与 Python 在真实路径上看到的是明文 —— 无需启动器、无需 shim、无需守护进程。",
    "",
    "##集成大幅简化",
    "业务二进制不被修改或重新封装。原生符号解析「开箱即用」，消除了 shim 方案必须绕开的一整类依赖顺序问题。",
    "",
    "##由内核访问门禁决定谁可解密",
    "只有授权进程（白名单/签名校验）才能通过挂载点拿到明文。拷贝类工具（cp、文件管理器、备份）拿到的是被剥除密钥的密文 —— 拷贝得到的是无用字节。",
    "",
    "##密钥位置",
    "嵌入在每个加密文件中。无挂载口令、无 keyring。代价：若直接拷贝 .enc/ 目录即可解密（见安全章节）。",
], body_size=13)

table_slide(prs, "设计对比一览", HDR,
    [
        ["层次", "仅用户态 —— 不改内核", "设备上需签名内核模块"],
        ["二进制", "每个被重新封装为 stub 启动器", "原封不动 —— 挂载而非封装"],
        ["解密路径", "守护进程解密到内存 memfd", "内核 FS 读取时解密（页缓存）"],
        ["注入", "每个进程注入 LD_PRELOAD shim", "无 —— 原生加载，真实路径"],
        ["/proc 视图", "memfd:<随机>（deleted）", "真实路径 —— 无 memfd 痕迹"],
        ["依赖处理", "需绕开符号顺序、dlopen、protobuf 去重", "原生 ld.so —— 问题消失"],
        ["密钥位置", "仅在启动器 + 守护进程二进制内", "嵌入每个文件"],
        ["磁盘被拷贝", "无密钥密文 —— 安全", "可解密 —— 需门禁 + 命名空间隔离"],
        ["现状", "交付中 / 生产", "候选；生产形态已端到端测试"],
    ], col_widths=COLW, font=11)

content_slide(prs, "AES 密钥在哪里？（安全核心结论）", [
    "##现行方案 —— 密钥集中在可执行文件中",
    "32 字节的 AES-256 密钥被追加在每个受保护 exe 的启动器内、以及守护进程二进制内。550 个库文件加密时不含密钥（无密钥容器）。",
    "客户把库文件夹拷到 U 盘，得到的是无密钥密文 —— 无法解密。密钥只随程序二进制本身传播。",
    "",
    "##kmod2 —— 密钥嵌入每个文件",
    "每个加密文件在尾部携带自己的密钥；内核在解密时读取它。没有单独的密钥文件、没有口令、没有需要管理的 keyring。",
    "权衡：直接拷贝 .enc/ 目录是可解密的。保护因此取决于：(1) 挂载层对未授权读者剥除密钥；(2) 访问门禁；(3) 把 .enc/ 目录置于操作员无法浏览的隔离挂载命名空间中。",
    "",
    "##两者共有的硬性限制",
    "CPU 执行的是明文，因此明文必然在内存中，拥有运行中设备 root 权限者始终能读取它。我方威胁模型将此排除在外：在设备旁的操作员是非技术人员，而有能力的竞争对手只能拿到被拷出的字节，永远接触不到运行中的设备。",
    "",
    "##对两者的推荐兜底：出厂时按单元用 TPM 封存密钥",
    "被封存的密钥永不离开设备、在不同硬件上无法解封 —— 使被拷出的密文彻底无用。",
], body_size=13)

# ════════════════ kmod2 图解 ════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, BG_DARK)
ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.55), Inches(2), Pt(4))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
tb = add_textbox(slide, Inches(1.2), Inches(2.4), Inches(7.8), Inches(1.2))
set_text(tb.text_frame, "kmod2 图解", size=40, color=WHITE, bold=True)
tb = add_textbox(slide, Inches(1.2), Inches(3.8), Inches(8.2), Inches(1.0))
set_text(tb.text_frame, "分层架构 · 解密读路径 · 解密门禁 · 磁盘容器格式", size=18, color=GRAY)

# 图 1 · 分层架构
slide = diagram_title(prs, "图 1 · kmod2 分层架构（解密内核文件系统）",
    "读取：穿过 antirevfs 即时解密；写入：落到 overlay 上层，绝不触碰下层密文")
dbox(slide, 0.45, 1.15, 5.4, 0.9, "① 业务应用 · glibc · ld.so · Python",
     "看到：真实路径 + 明文（与未加密完全一样）", border=GREEN, tcolor=GREEN)
darrow(slide, 3.0, 2.08, 0.3, 0.3); dlabel(slide, 3.4, 2.08, 1.7, "读取 ↓ 解密", ACCENT2)
dbox(slide, 0.45, 2.45, 5.4, 1.45, "② antirevfs 挂载点（真实安装路径）",
     "read_folio 读取时即时解密\n内核页缓存跨进程共享 · /proc 真实路径\n无 memfd · 无守护进程 · 无 LD_PRELOAD",
     border=ACCENT, tcolor=ACCENT)
darrow(slide, 3.0, 3.95, 0.3, 0.3); dlabel(slide, 3.4, 3.95, 1.9, "写入 ↓ overlay", ORANGE)
dbox(slide, 0.45, 4.3, 5.4, 0.8, "③ overlay 可写上层",
     "运行时写入（lock / log / pid）落这里，绝不触碰密文", border=ORANGE, tcolor=ORANGE)
darrow(slide, 3.0, 5.15, 0.3, 0.3)
dbox(slide, 0.45, 5.5, 5.4, 1.2, "④ 下层 .enc/ 密文树（磁盘上）",
     "ANTREV01 容器：magic+iv+tag+ct+key+magic\n（密钥内嵌每文件）", border=RED, tcolor=RED)
dbox(slide, 6.1, 1.15, 3.5, 1.7, "解密门禁 (gate.c)",
     "授权进程 → 明文（共享页缓存）\ncp / 备份 / 文件管理器 → 剥密钥密文 或 EACCES\nexec-load 按程序自身路径；数据读按调用进程身份",
     border=ORANGE, tcolor=ORANGE, tsize=12, bsize=10)
dbox(slide, 6.1, 3.0, 3.5, 1.35, "为什么更简单",
     "原生符号解析，无 LD_PRELOAD\n无守护、无 memfd 指纹\n依赖顺序 / protobuf 问题消失",
     border=GREEN, tcolor=GREEN, tsize=12, bsize=10)
dbox(slide, 6.1, 4.5, 3.5, 2.0, "威胁模型权衡",
     "密钥内嵌每文件 → 裸拷 .enc/ 可解密\n防护靠：门禁剥密钥 + 命名空间隔离\n出厂 TPM 封存密钥兜底",
     border=RED, tcolor=RED, tsize=12, bsize=10)

# 图 2 · 解密读路径
slide = diagram_title(prs, "图 2 · 解密读路径：read_folio 读取时解密",
    "首次缺页一次性整文件 GCM 解密并校验 tag，之后按页从页缓存服务；按页随机解密不可行")
_steps = [
    ("① 读 / 缺页", "程序读挂载下文件\n触发页错误", ACCENT),
    ("② read_folio", "内核读取时\n进入解密", ACCENT2),
    ("③ 取密钥", "从文件尾部\n读内嵌密钥\n用后清零", ORANGE),
    ("④ GCM 解密", "整文件一次性\n解密 + 校验 tag", GREEN),
    ("⑤ 页缓存", "写每-inode 缓冲\n→ 内核页缓存", ACCENT2),
    ("⑥ 返回明文", "服务页\n跨进程共享", GREEN),
]
_bx, _bw, _gap, _by, _bh = 0.35, 1.42, 0.16, 2.5, 1.95
for _i, (_t, _b, _c) in enumerate(_steps):
    _x = _bx + _i * (_bw + _gap)
    dbox(slide, _x, _by, _bw, _bh, _t, _b, border=_c, tcolor=_c, tsize=12, bsize=10)
    if _i < len(_steps) - 1:
        darrow(slide, _x + _bw + 0.01, _by + _bh / 2 - 0.11, _gap - 0.02, 0.22,
               color=GRAY, shape=MSO_SHAPE.RIGHT_ARROW)
dbox(slide, 0.35, 4.75, 9.27, 1.55, "未授权读者（门禁拦截）",
     "跳过解密 → 返回去掉 40B 密钥尾部的密文（cat 见 ANTREV01 magic，objdump 无法识别）。\n"
     "GCM 对整条消息做认证，无法按页随机解密，因此首次整文件解密、tag 只校验一次。",
     border=RED, tcolor=RED, tsize=13, bsize=11)

# 图 3 · 解密门禁决策
slide = diagram_title(prs, "图 3 · 解密门禁：每进程授权决策（仅作用于加密文件）",
    "exec-load 按『程序自身路径』授权；其它一切（库加载 / cp / source）按调用进程 current->mm->exe_file 授权")
dbox(slide, 3.1, 1.05, 3.8, 0.55, "打开挂载下的加密文件 (ANTREV01)", border=ACCENT, tsize=12)
darrow(slide, 4.85, 1.62, 0.3, 0.26)
dbox(slide, 4.0, 1.92, 2.0, 0.92, "exec-load ?", border=ACCENT2, tcolor=WHITE,
     tsize=13, shape=MSO_SHAPE.DIAMOND)
dconnect(slide, 4.0, 2.38, 2.05, 3.1); dlabel(slide, 2.3, 2.42, 1.6, "是 · 执行", ACCENT2, 10, True)
dbox(slide, 1.0, 3.1, 2.1, 0.95, "程序自身\n在白名单 ?", border=ACCENT2, tcolor=WHITE,
     tsize=12, shape=MSO_SHAPE.DIAMOND)
dconnect(slide, 1.55, 4.05, 1.15, 4.45); dlabel(slide, 0.2, 4.1, 1.0, "是", GREEN, 10, True)
dbox(slide, 0.35, 4.45, 1.6, 0.6, "✓ 允许执行", border=GREEN, tcolor=GREEN, tsize=12)
dconnect(slide, 2.55, 4.05, 3.0, 4.45); dlabel(slide, 3.0, 4.1, 1.0, "否", RED, 10, True)
dbox(slide, 2.15, 4.45, 1.75, 0.6, "✗ EACCES (126)", border=RED, tcolor=RED, tsize=12)
dconnect(slide, 6.0, 2.38, 7.2, 3.1); dlabel(slide, 6.05, 2.42, 2.0, "否 · 数据读", ACCENT2, 10, True)
dbox(slide, 6.0, 3.1, 2.4, 0.95, "调用进程\n在白名单 ?", border=ACCENT2, tcolor=WHITE,
     tsize=12, shape=MSO_SHAPE.DIAMOND)
dconnect(slide, 6.6, 4.05, 5.7, 4.45); dlabel(slide, 4.95, 4.1, 1.0, "是", GREEN, 10, True)
dbox(slide, 4.6, 4.45, 2.05, 0.78, "✓ 解密 → 明文\n（共享页缓存）", border=GREEN, tcolor=GREEN,
     tsize=12, bsize=10)
dconnect(slide, 7.9, 4.05, 8.5, 4.45); dlabel(slide, 8.5, 4.1, 1.0, "否", RED, 10, True)
dbox(slide, 7.6, 4.45, 2.0, 0.9, "gate_passthrough\n_cipher ?", border=ACCENT2, tcolor=WHITE,
     tsize=11, shape=MSO_SHAPE.DIAMOND)
dconnect(slide, 8.2, 5.35, 7.75, 5.7); dlabel(slide, 6.75, 5.42, 1.0, "是", ORANGE, 10, True)
dbox(slide, 6.9, 5.7, 1.7, 0.55, "剥密钥密文", border=ORANGE, tcolor=ORANGE, tsize=11)
dconnect(slide, 9.0, 5.35, 9.2, 5.7); dlabel(slide, 8.7, 5.42, 0.8, "否", RED, 10, True)
dbox(slide, 8.65, 5.7, 1.1, 0.55, "EACCES", border=RED, tcolor=RED, tsize=11)

# 图 4 · 磁盘容器格式
slide = diagram_title(prs, "图 4 · 磁盘容器格式：内嵌密钥尾部（无挂载密钥）",
    "授权读者：内核读尾部密钥即时解密（栈上拷贝用后清零）；未授权读者：尾部 40B 被剥除")
_segs = [("ANTREV01\n8B", 1.1, ACCENT), ("IV\n12B", 1.0, ACCENT2),
         ("TAG\n16B", 1.2, ORANGE), ("密文 ct\n（明文经 GCM 加密）", 3.0, GREEN),
         ("KEY\n32B", 1.5, RED), ("ANTREV01\n8B", 1.1, ACCENT)]
_sx, _sy, _sh = 0.55, 2.7, 1.05
_xs = []; _cx = _sx
for (_t, _w, _c) in _segs:
    _xs.append((_cx, _w, _c, _t)); _cx += _w
for (_x, _w, _c, _t) in _xs:
    dbox(slide, _x, _sy, _w, _sh, _t, border=_c, tcolor=_c, tsize=11)
dlabel(slide, _xs[0][0], 2.28, (_xs[2][0] + _xs[2][1]) - _xs[0][0],
       "HDR = 36B（magic+iv+tag）", ACCENT2, 11, True)
dlabel(slide, _xs[3][0], 2.28, _xs[3][1], "明文长度 = 文件 − 76B", GREEN, 11, True)
dlabel(slide, _xs[4][0], 2.28, (_xs[5][0] + _xs[5][1]) - _xs[4][0],
       "TRAILER = 40B（key+magic）", RED, 11, True)
dlabel(slide, _xs[4][0], _sy + _sh + 0.05, (_xs[5][0] + _xs[5][1]) - _xs[4][0],
       "未授权读者 → 去掉这 40B → 无密钥容器", RED, 10.5, True)
dbox(slide, 0.55, 4.55, 9.05, 1.75, "要点",
     "密钥随文件走：裸拷 .enc/ 可解密 → 必须命名空间隔离 + 出厂 TPM 封存兜底。\n"
     "授权读者：读尾部密钥即时解密。\n"
     "未授权读者：尾部 40B 被剥除，cp / vim / objdump 只得到无用的无密钥容器（objdump：无法识别）。",
     border=ACCENT, tcolor=ACCENT, tsize=13, bsize=11)

# ════════════════ 第 2 部分：CI 构建之后 ════════════════
section_slide(prs, 2, "CI 构建之后", "软件包包含什么 —— 新增文件与被修改文件")

content_slide(prs, "发布流水线 —— 两个阶段（两种方案通用）", [
    "##阶段 A —— 编译（CI，与密钥无关、可复用）",
    "从源码构建保护工具链。产出不含客户密钥的各架构二进制。同一份阶段 A 产物在所有部署间复用。",
    "  现行：stub 启动器、antirev_shim_<arch>.so、lrxd-<arch> 守护进程、打包工具",
    "  kmod2：antirevfs.ko（每个内核版本一个，经 DKMS）、挂载工具",
    "",
    "##阶段 B —— 打包（按部署、使用机密密钥）",
    "将业务安装目录 + 部署密钥 + 阶段 A 产物，打包成加密的可交付软件包。这一步注入/使用密钥并产出交付给客户的成果。",
    "  现行：antirev-pack.py → 重新封装的 exe + 加密库 + 守护进程",
    "  kmod2：antirev-fs-pack.py → 密文下层目录 .enc/ + manifest",
    "",
    "##为什么这个划分对发布管理很重要",
    "阶段 A 是普通 CI（构建一次、签名、归档）。阶段 B 是受控的、携带密钥的步骤，应在受保护的发布环境中执行。密钥是唯一机密；其余一切都可从源码复现。",
], body_size=14)

table_slide(prs, "CI 之后：软件包是什么样子", HDR,
    [
        ["新增二进制", "lrxd-<arch> 守护进程、antirev_shim_<arch>.so、设置 LD_PRELOAD 的启动脚本", "antirevfs.ko（签名）、挂载工具、开机/systemd 挂载单元"],
        ["新增数据文件", "—", ".enc/ 密文目录、manifest.json、proj-pack.yaml、/etc/authorized_apps.txt"],
        ["被修改", "每个 exe → stub+密文+密钥；每个 .so → 无密钥密文", "无 —— 安装路径变为 .enc/ 之上的挂载点"],
        ["Python 加载", "用 antirev_client 加载加密 .so", "不变 —— 通过挂载点原生加载"],
        ["密钥成果物", "单独的 64 位十六进制密钥文件（0600）/ TPM 封存", "嵌入每个文件 —— 无需单独发布密钥"],
        ["第三方库", "原封不动（明文共存）", "在目录中保留明文（已加入黑名单）"],
    ], col_widths=COLW, font=11)

# ════════════════ 第 3 部分：运行时 ════════════════
section_slide(prs, 3, "运行时", "安装目录与运行进程是什么样子")

table_slide(prs, "设备上的运行时形态", HDR,
    [
        ["磁盘上", "stub 封装的 exe + 无密钥密文库 + lrxd 守护进程", ".enc/ 下层目录（隐藏）+ 显示真实路径的挂载点"],
        ["额外进程", "一个长驻 lrxd 守护进程", "无"],
        ["注入库", "每个进程注入 LD_PRELOAD shim", "无 —— 应用原样运行"],
        ["/proc/<pid>/exe", "memfd:<随机>（deleted）", "真实路径"],
        ["/proc maps", "memfd 条目，无库路径", "真实库路径"],
        ["其他痕迹", "抽象 Unix socket、/tmp/antirev_* 软链目录、打开的 memfd", "lsmod 显示 antirevfs、挂载条目、共享页缓存、内核门禁"],
        ["可写 bin/lib", "普通文件", "在只读 antirevfs 之上叠加 overlay 上层"],
    ], col_widths=COLW, font=11)

content_slide(prs, "运行时：kmod2 的可写视图细节", [
    "##该挂载默认是只读的",
    "antirevfs 只提供已交付的加密内容。但真实业务软件会在其二进制旁边写入 lock/pid/log 文件（例如 GUI 把 QtApplication.pid 写入 bin/）。",
    "",
    "##解决：用标准 overlay 提供可写上层",
    "antirev-mount-rw 在解密文件系统之上叠加一个标准可写层。解密后的库读取下沉到 antirevfs；应用的运行时写入落在可写上层，绝不触及密文 .enc/ 目录。参考部署中 bin/ 与 lib/ 都以此方式挂载。",
    "",
    "##操作员可见的整体画面",
    "一个看起来正常、应用照常读写的 bin/ 与 lib/ —— 而磁盘上的真实字节保持加密，且只有授权进程能通过挂载点取得明文。",
], body_size=14)

# ── 开发者自验证场景 ────────────────────────────────────────────────
table_slide(prs, "开发者自验证：新编一个二进制的影响",
    ["维度", "memfd + 守护进程（现行）", "kmod2（antirevfs）"],
    [
        ["要测新二进制需做什么", "先打包（加密 + 封装 stub）→ 启动守护进程 → 运行", "加密进 .enc/ → 挂载 → 运行（也可全程明文开发，最后加密一次验证）"],
        ["与编译产物是否一致", "被重新封装，非字节一致；进程模型被改变", "解密后字节一致，完全透明"],
        ["调试体验（gdb / core）", "难 —— /proc 指向 memfd、readlink 被伪装、符号隐藏、core 指向 memfd", "易 —— 真实路径、原生符号，gdb / core dump 正常"],
        ["保护态专属 bug", "有一类只在保护下出现（符号顺序、protobuf 去重、隐式依赖 / NO_PRELOAD）", "基本没有 —— 原生加载，明文行为 == 加密行为"],
        ["权限要求", "用户态即可，无需 root", "需 root（insmod + 挂载）；开发期可 gate_enforce=0 关门禁"],
        ["迭代回路", "改 → 重新打包 → 重启守护进程 → 跑", "改 → 重新加密该文件 → 重新挂载 → 跑"],
        ["新增依赖库", "需更新 needed-libs 元数据 / 重打包 exe", "无需 —— 原生 DT_NEEDED 解析"],
        ["门禁对开发的影响", "无门禁概念", "若开门禁，需把 shell / 解释器 / gdb 按 basename 加白名单"],
        ["CI 中自验证", "纯用户态，普通 runner 即可跑", "需特权容器 / 内核模块，CI 更重"],
        ["加密对开发者可感知性", "高 —— 封装 / 守护进程 / memfd 痕迹外显", "低 —— 真实路径，看似普通目录"],
    ], col_widths=(Inches(1.95), Inches(3.7), Inches(3.75)), font=10)

content_slide(prs, "开发者自验证：在「不让人感知加密」前提下验证改动", [
    "##诉求：开发者改了自己的 exe/.so、放进文件系统验证，且其他开发者不应察觉公司软件包被加密",
    "",
    "##通用前提（两方案都适用）：开发者只在明文树上开发自测",
    "加密只发生在 CI 的打包阶段（阶段 B）。开发者在未加密的源码/构建树上开发、运行、调试 —— 与无保护时完全一致，天然无感知。这是首选，也是大多数自验证的答案。",
    "",
    "##当改动必须跑在「受保护包」内时，两方案差别很大：",
    "kmod2：真实路径透明 —— 把自编二进制丢进可写 overlay 即明文落盘，与解密后的库在真实路径上共存；开发者看到的就是普通目录，无 memfd / 无守护 / 无 LD_PRELOAD，察觉不到加密。",
    "  注意：门禁开启时，未白名单的自编二进制读加密库会得到 EACCES（像普通权限错误，反而可能引人怀疑）→ 给「开发盒」用 gate_enforce=0，或按 basename 临时白名单。",
    "memfd + 守护进程：自编二进制要用到加密库，必须经保护启动脚本（LD_PRELOAD + 守护 socket）；而 stub 封装、lrxd 守护进程、memfd 痕迹都外显 → 很难对开发者隐藏加密。实务上只能让开发者在明文树验证，保护态由 CI/专人验证。",
    "",
    "##结论",
    "「不让开发者感知加密」这一诉求：kmod2 在文件系统层天然满足（透明、真实路径）；memfd + 守护进程因机制外显，基本只能靠『开发者只碰明文树』来实现。",
], body_size=13)

content_slide(prs, "如何做到「不感知加密」：两套配方", [
    "##kmod2 配方（保护态本身透明，配置到位即可）",
    "把包以 antirevfs + overlay-rw + passdata 挂在真实安装路径；开发盒上 gate_enforce=0。",
    "自替换：cp 自编二进制进挂载点 → 落到 overlay 上层（明文、原样 served、不过门禁）；同名文件自动覆盖下层解密库，他的版本即时生效。",
    "自验证：正常启动即可 —— 真实路径，无 memfd / 无守护 / 无 LD_PRELOAD，行为 == 明文。",
    "",
    "##memfd + 守护进程 配方（机制外显，只能让开发者待在明文面）",
    "主路径：开发者只在未加密的构建树上替换 + 验证 —— 世界里没有加密，天然无感知；加密只在 CI 阶段 B 做，保护态由 CI/专人验。",
    "代价：验不到「明文能跑、加密就崩」那类 bug（符号顺序 / protobuf / 隐式依赖 / NO_PRELOAD），由 CI 兜底。",
    "若必须在保护态跑：用 wrapper 脚本把 LD_PRELOAD / 守护 socket 等环境藏进去 → 藏掉命令，但 /proc memfd、ps lrxd 仍外显；且同名加密库的明文替换会被守护的符号目录盖掉。",
], body_size=13)

table_slide(prs, "不感知加密：两架构对照", HDR,
    [
        ["隐藏发生在哪", "藏不住运行态，只能让开发者待在明文树", "藏在内核读路径，保护态本身透明"],
        ["自替换动作", "明文树里 cp；保护态同名加密库替换会被守护盖掉", "挂载点 cp → overlay 上层明文，同名自动覆盖、即时生效"],
        ["自验证范围", "只验功能；保护态专属 bug 要 CI/专人验", "在真实路径上直接验，行为 == 明文"],
        ["必做的一步", "wrapper 脚本藏环境（仍藏不掉 memfd / 守护）", "开发盒 gate_enforce=0（或 basename 白名单）"],
        ["残余暴露", "/proc memfd、ps lrxd、LD_PRELOAD 外显", "仅 lsmod / mount 可见模块；正常流程无感知"],
    ], col_widths=COLW, font=11)

# ════════════════ 第 4 部分：部分升级 ════════════════
section_slide(prs, 4, "部分升级", "现场只替换 / 新增部分库")

table_slide(prs, "部分升级 —— 替换 / 新增部分库",
    ["步骤", "memfd + 守护进程（现行）", "kmod2（antirevfs）"],
    [
        ["密钥要求", "必须复用同一部署密钥（所有文件 + 守护进程共用）", "文件自带密钥；为一致性复用项目密钥"],
        ["重新加密", "用现有密钥加密变更/新增的 .so", "把变更/新增 ELF 加密进 .enc/（内嵌密钥形式）"],
        ["放置文件", "把密文库放入目录", "在 .enc/ 下层目录中替换文件"],
        ["新增依赖边", "重新打包受影响的 exe（needed-libs 列表变化）", "无需 —— 原生解析"],
        ["生效", "重启 lrxd 守护进程（重新扫描并解密）", "重新挂载该目录（分类按文件缓存）"],
        ["注意点", "必须重启守护进程", "需要静止的业务栈 —— 已映射的库会钉住模块"],
    ], col_widths=COLW, font=11)

content_slide(prs, "部分升级 —— 给发布规划的要点", [
    "##密钥是每次升级都必须保留的唯一成果物",
    "现行方案：强制 —— 部署中每个文件都绑定同一密钥；新库必须用它加密，否则守护进程无法提供一致的集合。",
    "kmod2：约束更松 —— 文件自描述其密钥 —— 但保留单一项目密钥文件可让 manifest 与工具链保持整洁。",
    "",
    "##实务建议",
    "把每个部署的密钥存入发布保险库，并在出厂时按单元用 TPM 封存一份。后续补丁就无需发布密钥 —— 技术员用保留的密钥重新加密并放入文件即可。",
    "",
    "##每种方案各自的一个操作坑",
    "现行：别忘了重启守护进程（它在启动时解密）。",
    "kmod2：别忘了重新挂载（inode 分类有缓存）；并需要静止的业务栈，因为已映射的库会钉住模块。",
], body_size=14)

# ════════════════ 第 5 部分：完整升级 ════════════════
section_slide(prs, 5, "完整升级", "替换整套受保护软件")

table_slide(prs, "完整升级 —— 替换整套受保护软件",
    ["步骤", "memfd + 守护进程（现行）", "kmod2（antirevfs）"],
    [
        ["停止", "业务栈 + lrxd 守护进程", "业务栈（已映射的 .so 会钉住模块）"],
        ["拆除", "—", "卸下挂载（antirev-mount-rw --down）"],
        ["密钥选择", "保留（原地替换）或轮换（整体重打包、重新 TPM 封存）", "保留或轮换；每个文件自带密钥"],
        ["部署", "替换打包目录（exe、库、守护进程、shim、脚本）", "替换 .enc/ 下层目录 + manifest"],
        ["内核依赖", "无 —— 纯用户态替换", ".ko 须匹配内核；DKMS 重建 + Secure Boot 签名"],
        ["生效", "先启动守护进程，再启动业务栈", "重新挂载（antirev-remount-proj.sh），启动业务栈"],
        ["回滚", "保留旧打包目录 + 密钥；换回", "保留旧 .enc/ 目录；重新挂载"],
    ], col_widths=COLW, font=11)

content_slide(prs, "完整升级 —— 发布管理须把控什么", [
    "##现行方案 —— 完全自包含、无主机耦合",
    "完整升级就是一次用户态文件替换加守护进程重启。不依赖客户内核。在异构客户硬件上最易支持。",
    "",
    "##kmod2 —— 多了一个需把控的内核模块生命周期",
    ".ko 必须为每个客户内核版本构建并签名。DKMS 会在内核升级时自动重建，但 Secure Boot 要求我们的模块签名密钥被预置信任。这是新的发布/QA 面：客户设备上的内核更新可能需要重建模块。",
    "",
    "##对两种路径的建议",
    "把完整升级作为带版本、已签名的捆绑包发布（阶段 A 产物 + 阶段 B 打包目录），并配有经过测试的回滚（保留旧目录 + 密钥）。把部署密钥当作发布机密管理，以出厂 TPM 封存作为现场兜底。",
], body_size=14)

# ════════════════ 第 6 部分：安全态势 ════════════════
section_slide(prs, 6, "安全态势、权衡与建议", "项目决策所依赖的全部要点")

content_slide(prs, "威胁模型（已达成共识）—— 始终牢记", [
    "##设备旁的对手 = 客户 / 操作员",
    "物理上拥有硬件，但是非技术人员：只会点 GUI 按钮、拖拽、拷到 U 盘。没有 root、挂载、解密或内存取证的概念。攻击 =「把软件文件夹拷下来交给别人」。",
    "",
    "##有能力的对手 = 竞争对手",
    "专业逆向工程师，但与运行中的设备解耦 —— 他们只能拿到操作员拷出的字节，永远拿不到运行中设备的 root。",
    "",
    "##因此整个问题归结为：",
    "「非技术用户能通过 GUI/U 盘拷走的，必须是密文。」",
    "在运行设备上的高级 root 攻击确实不在范围内：在设备旁的人做不到，能做到的人不在设备旁。",
], body_size=15)

table_slide(prs, "针对该威胁模型的安全打分", HDR,
    [
        ["拷贝库文件夹", "无密钥密文 —— 密钥不在库中（强）", "经挂载：被剥除密钥的密文（无用）"],
        ["拷贝整个磁盘目录", "无密钥密文 —— 安全", ".enc/ 可解密 —— 须置于隔离挂载命名空间"],
        ["明文暴露", "内存 memfd；memfd 痕迹暗示「受保护」", "内核页缓存；真实路径，无痕迹"],
        ["残余风险", "密钥集中在 exe + 守护进程二进制", "嵌入每个文件的密钥是弱点"],
        ["加固路径", "混淆 + TPM 封存", "包裹 / TPM 封存每文件密钥（待办）"],
        ["成熟度", "已在完整规模上生产验证", "候选；生产栈形态已端到端测试"],
    ], col_widths=COLW, font=11)

content_slide(prs, "建议与后续步骤", [
    "##把现行 memfd + 守护进程方案保留为交付路径",
    "它经过生产验证、完全用户态（不与客户内核耦合），且静态态势强：部署的主体（550 个库）是无密钥密文。",
    "",
    "##把 kmod2 作为战略候选继续推进",
    "它集成大幅简化（无 shim、无依赖顺序绕弯、二进制原封不动），且不留 memfd 指纹。真实生产部署形态 —— 进程管理器栈、访问门禁、可写视图 —— 现已有自动化端到端测试覆盖（x86-64 与 ARM64-under-qemu）。其主要缺口是内嵌密钥 + 内核模块生命周期。",
    "",
    "##要让 kmod2 达到可发布，须收尾以下事项：",
    "  1. 加固内嵌密钥 —— 混淆 / TPM 封存每文件尾部，使被拷出的 .enc/ 目录同样无法解密",
    "  2. 为 Secure Boot 做模块签名 + DKMS .deb/.rpm 打包",
    "  3. 把访问门禁从白名单（演示）升级为基于签名的授权",
    "  4. 验证面向 ARM64-on-x86 RTOS 从机的 Docker/容器命名空间路径",
    "",
    "##对两者：出厂按单元用 TPM 封存密钥",
    "这是单点杠杆最高的控制 —— 使任何被拷出的密文在脱离原硬件后即无用。",
], body_size=13)

# 收尾
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.7), Inches(2), Pt(4))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
tb = add_textbox(slide, Inches(1.5), Inches(1.2), Inches(7.5), Inches(1.4))
set_text(tb.text_frame, "一句话总结", size=34, color=WHITE, bold=True)
tb = add_textbox(slide, Inches(1.5), Inches(3.0), Inches(7.9), Inches(4.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "现行方案 = 已验证、用户态、密钥藏于可执行文件、库为无密钥密文 —— 继续交付。",
         size=18, color=LIGHT)
add_para(tf, "", size=8)
add_para(tf, "kmod2 = 更简洁干净、密钥按文件内嵌 —— 在密钥加固且内核模块完成签名/打包后大有可为。",
         size=18, color=LIGHT)
add_para(tf, "", size=8)
add_para(tf, "对两者而言：部署密钥是核心资产 —— 在发布侧存入保险库，在出厂时按单元用 TPM 封存。",
         size=18, color=ACCENT2, bold=True)

# ── 保存 ─────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "kmod2_vs_daemon_report_zh.pptx")
prs.save(out)
print("已保存:", out, "—", len(prs.slides._sldIdLst), "页")
